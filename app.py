from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify, send_file
from werkzeug.security import generate_password_hash, check_password_hash
import sqlite3
from datetime import datetime, timedelta
import json
from functools import wraps
import os
import csv
import io
from collections import defaultdict
from pptx import Presentation
import pandas as pd
import traceback

app = Flask(__name__)
app.secret_key = 'your-secret-key-here-change-in-production'

# Configuration
DATABASE = 'qualite_ADI_final.sqlite'
UPLOAD_FOLDER = 'static/uploads'

def get_db_connection():
    """Créer une connexion à la base de données"""
    conn = sqlite3.connect(DATABASE)
    conn.row_factory = sqlite3.Row
    return conn

def init_database():
    """Initialiser la base de données avec les tables nécessaires"""
    conn = get_db_connection()
    
    # Table des utilisateurs avec protection et traçabilité
    conn.executescript('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            nom TEXT NOT NULL,
            prenom TEXT NOT NULL,
            email TEXT UNIQUE NOT NULL,
            nom_compte TEXT UNIQUE NOT NULL,
            mot_de_passe TEXT NOT NULL,
            role TEXT NOT NULL CHECK (role IN ('Administration', 'Chargé Qualité/Animateur Qualité', 'Autre')),
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            cree_par TEXT,
            modifie_par TEXT,
            is_protected INTEGER DEFAULT 0
        );
        
        CREATE TABLE IF NOT EXISTS seuils (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            atelier TEXT NOT NULL,
            indicateur TEXT NOT NULL,
            valeur_seuil REAL NOT NULL,
            modifie_by TEXT NOT NULL,
            date_modification TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS pieces_retouchees (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            quantite INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS pieces_rebutees (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            quantite INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS pieces_controlees (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            quantite INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS cout_rebut (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            montant REAL NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS cnq (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            valeur REAL NOT NULL,
            seuil REAL NOT NULL,
            valeur_ytd DECIMAL(10, 3) NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS pieces_reclamees (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            quantite INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS pieces_exportees (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            quantite INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS quantite_alertes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            nombre_alertes INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        
        CREATE TABLE IF NOT EXISTS reclamations_officielles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mois INTEGER NOT NULL CHECK (mois BETWEEN 1 AND 12),
            annee INTEGER NOT NULL,
            atelier TEXT NOT NULL,
            uap TEXT NOT NULL,
            nombre INTEGER NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
    ''')
    
    # CORRECTION 1: Vérifier l'admin existant avec protection contre les erreurs
    try:
        admin_exists = conn.execute(
            "SELECT COUNT(*) FROM users WHERE nom_compte = ?", 
            ('skander.chniti',)
        ).fetchone()[0]
        
        if admin_exists == 0:
            # CORRECTION 2: Import de generate_password_hash à l'intérieur de la fonction
            from werkzeug.security import generate_password_hash
            
            hashed_password = generate_password_hash('@@skanderchniti1$$')
            
            # CORRECTION 3: Nom de compte cohérent (minuscules)
            conn.execute('''
                INSERT INTO users (nom, prenom, email, nom_compte, mot_de_passe, role, cree_par, is_protected)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                'Chniti', 
                'Skander', 
                'skander.chniti@hutchinson.com', 
                'skander.chniti',  # Cohérent avec la vérification (minuscules)
                hashed_password, 
                'Administration',
                'Système', 
                1  # Protégé par défaut
            ))
            print("✅ Utilisateur admin créé et protégé")
        else:
            # L'admin existe déjà, s'assurer qu'il est protégé
            conn.execute(
                "UPDATE users SET is_protected = 1 WHERE nom_compte = ?",
                ('skander.chniti',)
            )
            print("✅ Compte admin existant mis à jour et protégé")
            
    except Exception as e:
        print(f"⚠️  Erreur lors de la gestion de l'admin: {e}")
    
    # CORRECTION 4: Vérifier que la fonction init_default_thresholds existe
    try:
        # Initialiser les seuils par défaut si ils n'existent pas
        init_default_thresholds(conn)
        print("✅ Seuils par défaut initialisés")
    except NameError:
        print("ℹ️  Fonction init_default_thresholds non trouvée, ignorée")
        # Vous pouvez ajouter l'initialisation des seuils manuellement ici si besoin
        pass
    except Exception as e:
        print(f"⚠️  Erreur lors de l'initialisation des seuils: {e}")
    
    conn.commit()
    conn.close()
    print("✅ Base de données initialisée avec succès")

# FONCTION MANQUANTE: Ajouter cette fonction si elle n'existe pas
def init_default_thresholds(conn):
    """Initialiser les seuils par défaut pour les indicateurs"""
    
    # Vérifier si des seuils existent déjà
    existing_thresholds = conn.execute("SELECT COUNT(*) FROM seuils").fetchone()[0]
    
    if existing_thresholds == 0:
        # Seuils par défaut pour les différents indicateurs
        default_thresholds = [
            ('Tous', 'Taux de Rebut (%)', 2.0),
            ('Tous', 'Taux de Retouche (%)', 3.0),
            ('Tous', 'Coût de Rebut (€)', 1000.0),
            ('Tous', 'CNQ', 95.0),
            ('Tous', 'Taux de Réclamation (%)', 1.0),
            ('Tous', 'Nombre d\'Alertes', 5),
            ('Tous', 'Réclamations Officielles', 2),
        ]
        
        for atelier, indicateur, valeur in default_thresholds:
            conn.execute('''
                INSERT INTO seuils (atelier, indicateur, valeur_seuil, modifie_by)
                VALUES (?, ?, ?, ?)
            ''', (atelier, indicateur, valeur, 'Système'))
        
        print("✅ Seuils par défaut créés")
    else:
        print("ℹ️  Seuils déjà présents dans la base")

# ALTERNATIVE: Si vous voulez une version plus robuste
def init_database_robust():
    """Version robuste avec gestion d'erreurs complète"""
    try:
        conn = get_db_connection()
        print("🔧 Connexion à la base de données établie")
        
        # Création des tables
        print("📋 Création des tables...")
        conn.executescript('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                nom TEXT NOT NULL,
                prenom TEXT NOT NULL,
                email TEXT UNIQUE NOT NULL,
                nom_compte TEXT UNIQUE NOT NULL,
                mot_de_passe TEXT NOT NULL,
                role TEXT NOT NULL CHECK (role IN ('Administration', 'Chargé Qualité/Animateur Qualité', 'Autre')),
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                cree_par TEXT,
                modifie_par TEXT,
                is_protected INTEGER DEFAULT 0
            );
            
            -- Autres tables... (copiez vos autres CREATE TABLE ici)
        ''')
        
        # Gestion de l'utilisateur admin
        print("👤 Vérification de l'utilisateur admin...")
        admin_check = conn.execute(
            "SELECT id, is_protected FROM users WHERE nom_compte = ?", 
            ('skander.chniti',)
        ).fetchone()
        
        if admin_check:
            # Admin existe, vérifier la protection
            if admin_check['is_protected'] != 1:
                conn.execute(
                    "UPDATE users SET is_protected = 1 WHERE id = ?",
                    (admin_check['id'],)
                )
                print("🔒 Protection activée pour l'admin existant")
            else:
                print("✅ Admin déjà protégé")
        else:
            # Créer l'admin
            from werkzeug.security import generate_password_hash
            hashed_password = generate_password_hash('@@skanderchniti1$$')
            
            conn.execute('''
                INSERT INTO users (nom, prenom, email, nom_compte, mot_de_passe, role, cree_par, is_protected)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                'Chniti', 'Skander', 'skander.chniti@hutchinson.com', 
                'skander.chniti', hashed_password, 'Administration', 'Système', 1
            ))
            print("✅ Utilisateur admin créé et protégé")
        
        # Initialiser les seuils
        print("📊 Initialisation des seuils...")
        init_default_thresholds(conn)
        
        conn.commit()
        conn.close()
        print("🎉 Initialisation de la base de données terminée avec succès !")
        
    except sqlite3.Error as e:
        print(f"❌ Erreur SQLite: {e}")
        if 'conn' in locals():
            conn.close()
        raise
    except Exception as e:
        print(f"❌ Erreur générale: {e}")
        if 'conn' in locals():
            conn.close()
        raise



def init_default_thresholds(conn):
    """Initialiser les seuils par défaut pour tous les ateliers"""
    ateliers = [
        'Manchon', 'Collier de Fixation', 'Rack', 'Moulage',
        'Isolation Thermique', 'Isolation Souple', 'Composite',
        'Atelier de Visualisation'
    ]
    
    indicateurs = {
        'PPM Officiel': 100,
        'PPM Non Officiel': 500,
        'Taux de Rebut': 2.0,
        'Taux de Retouche': 5.0,
        'CNQ': 1000,
        'Coût de Rebut': 500,
        'Nombre de reclamation': 10 # j ai ajouter cela 
    }
    
    for atelier in ateliers:
        for indicateur, valeur_defaut in indicateurs.items():
            # Vérifier si le seuil existe déjà
            existing = conn.execute('''
                SELECT COUNT(*) FROM seuils 
                WHERE atelier = ? AND indicateur = ?
            ''', (atelier, indicateur)).fetchone()[0]
            
            if existing == 0:
                conn.execute('''
                    INSERT INTO seuils (atelier, indicateur, valeur_seuil, modifie_by)
                    VALUES (?, ?, ?, ?)
                ''', (atelier, indicateur, valeur_defaut, 'Système'))

                #*************************************************#

#login_required empêche les gens non connectés d'accéder à certaines pages. Il redirige vers la page de login s'ils ne sont pas connectés.
def login_required(f):
    """Décorateur pour vérifier l'authentification"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


#protéger certaines pages pour que seuls les administrateurs y accèdent.
def admin_required(f):
    """Décorateur pour vérifier les droits admin"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_role' not in session or session['user_role'] != 'Administration':
            flash('Accès non autorisé', 'error')
            return redirect(url_for('dashboard'))
        return f(*args, **kwargs)
    return decorated_function


#appeller la page d'accueil (/) via l’URL http://localhost:5000/
@app.route('/')
def home():
    """Page d'accueil"""
    return render_template('home.html')


#gère la connexion des utilisateurs en affichant le formulaire de login (GET) et en vérifiant les identifiants soumis (POST) pour les connecter ou afficher un message d’erreur.
@app.route('/login', methods=['GET', 'POST'])
def login():
    """Page de connexion"""
    if request.method == 'POST':
        nom_compte = request.form['nom_compte']
        mot_de_passe = request.form['mot_de_passe']
        
        conn = get_db_connection()
        user = conn.execute(
            'SELECT * FROM users WHERE nom_compte = ?', (nom_compte,)
        ).fetchone()
        conn.close()
        
        if user and check_password_hash(user['mot_de_passe'], mot_de_passe):
            session['user_id'] = user['id']
            session['user_name'] = f"{user['prenom']} {user['nom']}"
            session['user_role'] = user['role']
            session['user_email'] = user['email']
            
            flash(f'Bienvenue {session["user_name"]}!', 'success')
            return redirect(url_for('dashboard'))
        else:
            flash('Nom de compte ou mot de passe incorrect', 'error')
    
    return render_template('login.html')

#déconnecte l’utilisateur en effaçant sa session, lui affiche un message de succès, puis le renvoie à la page d’accueil.
@app.route('/logout')
@login_required
def logout():
    """Déconnexion"""
    session.clear()
    flash('Vous avez été déconnecté avec succès', 'info')
    return redirect(url_for('home'))


@app.route('/dashboard')
@login_required
def dashboard():
    """Dashboard principal selon le rôle"""
    role = session.get('user_role')
    
    if role == 'Administration':
        return render_template('admin/dashboard.html')
    elif role == 'Chargé Qualité/Animateur Qualité':
        return render_template('quality_manager/dashboard.html')
    else:  # Autre
        return render_template('viewer/dashboard.html')


# 1...........Routes pour l'administration........

# Routes complètes pour la gestion des utilisateurs avec protection admin

# Fonction utilitaire pour vérifier la protection d'un utilisateur
def is_user_protected(user_id):
    """Vérifier si un utilisateur est protégé"""
    conn = get_db_connection()
    result = conn.execute(
        'SELECT is_protected FROM users WHERE id = ?', 
        (user_id,)
    ).fetchone()
    conn.close()
    return result and result['is_protected'] == 1

# Route COMPLÈTEMENT MODIFIÉE pour lister les utilisateurs
@app.route('/admin/users')
@login_required
@admin_required
def admin_users():
    """Gestion des utilisateurs avec informations de traçabilité"""
    conn = get_db_connection()
    users = conn.execute('''
        SELECT *, 
               CASE WHEN is_protected = 1 THEN 'Oui' ELSE 'Non' END as protege,
               CASE 
                   WHEN updated_at > created_at THEN strftime('%d/%m/%Y', updated_at)
                   ELSE NULL 
               END as derniere_modification
        FROM users 
        ORDER BY 
            is_protected DESC,  -- Utilisateurs protégés en premier
            created_at DESC
    ''').fetchall()
    conn.close()
    return render_template('admin/users.html', users=users)

# Route MODIFIÉE pour créer un utilisateur
@app.route('/admin/create_user', methods=['POST'])
@login_required
@admin_required
def create_user():
    """Créer un nouvel utilisateur avec traçabilité"""
    data = request.form
    
    # Validation des champs requis
    required_fields = ['nom', 'prenom', 'email', 'nom_compte', 'mot_de_passe', 'role']
    for field in required_fields:
        if not data.get(field):
            flash(f'Le champ {field} est requis', 'error')
            return redirect(url_for('admin_users'))
    
    # Validation du mot de passe
    password = data['mot_de_passe']
    if len(password) < 8:
        flash('Le mot de passe doit contenir au moins 8 caractères', 'error')
        return redirect(url_for('admin_users'))
    
    try:
        conn = get_db_connection()
        hashed_password = generate_password_hash(password)
        
        # Nettoyer le nom de compte (enlever espaces, mettre en minuscules)
        nom_compte_clean = data['nom_compte'].strip().lower()
        
        conn.execute('''
            INSERT INTO users (nom, prenom, email, nom_compte, mot_de_passe, role, cree_par, is_protected)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            data['nom'].strip().title(),  # Capitaliser le nom
            data['prenom'].strip().title(),  # Capitaliser le prénom
            data['email'].strip().lower(),  # Email en minuscules
            nom_compte_clean, 
            hashed_password, 
            data['role'], 
            session['user_name'],  # Qui a créé l'utilisateur
            0  # Pas protégé par défaut
        ))
        
        conn.commit()
        conn.close()
        
        flash(f'Utilisateur {data["prenom"]} {data["nom"]} créé avec succès', 'success')
        
    except sqlite3.IntegrityError as e:
        if 'email' in str(e):
            flash('Cette adresse email est déjà utilisée', 'error')
        elif 'nom_compte' in str(e):
            flash('Ce nom de compte est déjà utilisé', 'error')
        else:
            flash('Nom de compte ou email déjà existant', 'error')
    except Exception as e:
        flash(f'Erreur lors de la création: {str(e)}', 'error')
    
    return redirect(url_for('admin_users'))

# Route MODIFIÉE pour modifier un utilisateur
@app.route('/admin/edit_user/<int:user_id>', methods=['POST'])
@login_required
@admin_required
def edit_user(user_id):
    """Modifier un utilisateur avec vérifications de sécurité"""
    
    conn = get_db_connection()
    
    # Récupérer les informations de l'utilisateur à modifier
    user_to_edit = conn.execute(
        'SELECT is_protected, nom_compte, nom, prenom FROM users WHERE id = ?', 
        (user_id,)
    ).fetchone()
    
    if not user_to_edit:
        flash('Utilisateur introuvable', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    # PROTECTION 1: Vérifier si l'utilisateur est protégé
    if user_to_edit['is_protected'] == 1:
        flash(f'Impossible de modifier le compte {user_to_edit["nom_compte"]} : compte protégé par le système', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    # PROTECTION 2: Empêcher l'auto-modification (optionnel, vous pouvez l'enlever)
    if user_id == session.get('user_id'):
        flash('Vous ne pouvez pas modifier votre propre compte depuis cette interface', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    data = request.form
    
    try:
        # Si un nouveau mot de passe est fourni, le valider
        if data.get('mot_de_passe'):
            password = data['mot_de_passe']
            if len(password) < 8:
                flash('Le nouveau mot de passe doit contenir au moins 8 caractères', 'error')
                conn.close()
                return redirect(url_for('admin_users'))
            
            hashed_password = generate_password_hash(password)
            conn.execute('''
                UPDATE users 
                SET nom = ?, prenom = ?, email = ?, role = ?, mot_de_passe = ?, 
                    updated_at = CURRENT_TIMESTAMP, modifie_par = ?
                WHERE id = ?
            ''', (
                data['nom'].strip().title(), 
                data['prenom'].strip().title(), 
                data['email'].strip().lower(), 
                data['role'], 
                hashed_password, 
                session['user_name'], 
                user_id
            ))
        else:
            # Modification sans changement de mot de passe
            conn.execute('''
                UPDATE users 
                SET nom = ?, prenom = ?, email = ?, role = ?, 
                    updated_at = CURRENT_TIMESTAMP, modifie_par = ?
                WHERE id = ?
            ''', (
                data['nom'].strip().title(), 
                data['prenom'].strip().title(), 
                data['email'].strip().lower(), 
                data['role'], 
                session['user_name'], 
                user_id
            ))
        
        conn.commit()
        flash(f'Utilisateur {data["prenom"]} {data["nom"]} modifié avec succès', 'success')
        
    except sqlite3.IntegrityError as e:
        if 'email' in str(e):
            flash('Cette adresse email est déjà utilisée par un autre utilisateur', 'error')
        else:
            flash('Erreur : données en conflit avec un utilisateur existant', 'error')
    except Exception as e:
        flash(f'Erreur lors de la modification: {str(e)}', 'error')
    finally:
        conn.close()
    
    return redirect(url_for('admin_users'))

# Route MODIFIÉE pour supprimer un utilisateur
@app.route('/admin/delete_user/<int:user_id>', methods=['POST'])
@login_required
@admin_required
def delete_user(user_id):
    """Supprimer un utilisateur avec vérifications de sécurité renforcées"""
    
    # PROTECTION 1: Empêcher l'auto-suppression
    if user_id == session.get('user_id'):
        flash('Vous ne pouvez pas supprimer votre propre compte', 'error')
        return redirect(url_for('admin_users'))
    
    conn = get_db_connection()
    user = conn.execute(
        'SELECT * FROM users WHERE id = ?', 
        (user_id,)
    ).fetchone()
    
    if not user:
        flash('Utilisateur introuvable', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    # PROTECTION 2: Vérifier si l'utilisateur est protégé
    if user['is_protected'] == 1:
        flash(f'Impossible de supprimer le compte {user["nom_compte"]} : compte protégé par le système', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    # PROTECTION 3: Vérifier s'il reste au moins un admin après suppression
    if user['role'] == 'Administration':
        admin_count = conn.execute(
            "SELECT COUNT(*) as count FROM users WHERE role = 'Administration' AND id != ?", 
            (user_id,)
        ).fetchone()['count']
        
        if admin_count == 0:
            flash('Impossible de supprimer ce dernier administrateur. Il doit y avoir au moins un administrateur dans le système.', 'error')
            conn.close()
            return redirect(url_for('admin_users'))
    
    # Procéder à la suppression
    try:
        conn.execute('DELETE FROM users WHERE id = ?', (user_id,))
        conn.commit()
        flash(f'Utilisateur {user["prenom"]} {user["nom"]} supprimé avec succès', 'success')
    except Exception as e:
        flash(f'Erreur lors de la suppression: {str(e)}', 'error')
    finally:
        conn.close()
    
    return redirect(url_for('admin_users'))

# NOUVELLE ROUTE pour vérifier les permissions
@app.route('/api/user_permissions/<int:user_id>')
@login_required
@admin_required
def check_user_permissions(user_id):
    """Vérifier les permissions pour un utilisateur"""
    
    conn = get_db_connection()
    user = conn.execute(
        'SELECT is_protected, nom_compte, role FROM users WHERE id = ?', 
        (user_id,)
    ).fetchone()
    
    if not user:
        conn.close()
        return jsonify({
            'can_modify': False,
            'can_delete': False,
            'message': 'Utilisateur introuvable'
        })
    
    # Vérifications des permissions
    can_modify = True
    can_delete = True
    message = "OK"
    
    # L'utilisateur ne peut pas se modifier/supprimer lui-même
    if user_id == session.get('user_id'):
        can_modify = False
        can_delete = False
        message = "Vous ne pouvez pas modifier votre propre compte"
    
    # Vérifier si l'utilisateur est protégé
    elif user['is_protected'] == 1:
        can_modify = False
        can_delete = False
        message = f"Le compte {user['nom_compte']} est protégé et ne peut pas être modifié"
    
    # Vérifier s'il s'agit du dernier admin
    elif user['role'] == 'Administration':
        admin_count = conn.execute(
            "SELECT COUNT(*) as count FROM users WHERE role = 'Administration' AND id != ?", 
            (user_id,)
        ).fetchone()['count']
        
        if admin_count == 0:
            can_delete = False
            message = "Impossible de supprimer le dernier administrateur du système"
    
    conn.close()
    
    return jsonify({
        'can_modify': can_modify,
        'can_delete': can_delete,
        'message': message,
        'is_protected': user['is_protected'] == 1
    })

# NOUVELLE ROUTE pour obtenir les statistiques des utilisateurs
@app.route('/api/users_stats')
@login_required
@admin_required
def get_users_stats():
    """Obtenir les statistiques des utilisateurs"""
    
    conn = get_db_connection()
    
    # Compter par rôle
    stats_roles = conn.execute('''
        SELECT role, COUNT(*) as count 
        FROM users 
        GROUP BY role
    ''').fetchall()
    
    # Compter les utilisateurs protégés
    protected_count = conn.execute(
        'SELECT COUNT(*) as count FROM users WHERE is_protected = 1'
    ).fetchone()['count']
    
    # Utilisateurs créés cette semaine
    week_ago = (datetime.now() - timedelta(days=7)).isoformat()
    recent_count = conn.execute(
        'SELECT COUNT(*) as count FROM users WHERE created_at > ?',
        (week_ago,)
    ).fetchone()['count']
    
    conn.close()
    
    return jsonify({
        'total_users': sum(stat['count'] for stat in stats_roles),
        'by_role': {stat['role']: stat['count'] for stat in stats_roles},
        'protected_count': protected_count,
        'recent_count': recent_count
    })

# FONCTION UTILITAIRE pour l'audit des actions utilisateurs
def log_user_action(action, target_user_id, details=""):
    """Enregistrer les actions sur les utilisateurs pour audit (optionnel)"""
    # Cette fonction pourrait être étendue pour créer une table d'audit
    # Pour l'instant, on peut juste logger dans les fichiers
    
    current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    current_user = session.get('user_name', 'Inconnu')
    
    log_message = f"[{current_time}] {current_user} - {action} - User ID: {target_user_id} - {details}"
    
    # Écrire dans un fichier de log (optionnel)
    try:
        with open('user_actions.log', 'a', encoding='utf-8') as f:
            f.write(log_message + '\n')
    except:
        pass  # Ignorer les erreurs de logging

# ROUTE MODIFIÉE avec logging
@app.route('/admin/create_user_with_log', methods=['POST'])
@login_required
@admin_required
def create_user_with_log():
    """Version avec logging de la création d'utilisateur"""
    
    result = create_user()  # Appeler la fonction existante
    
    # Logger l'action si succès
    data = request.form
    if "créé avec succès" in str(get_flashed_messages()):
        log_user_action("CREATION", 0, f"Nouvel utilisateur: {data.get('prenom')} {data.get('nom')} ({data.get('role')})")
    
    return result

# NOUVELLE ROUTE pour changer le statut de protection (pour super admin uniquement)
@app.route('/admin/toggle_protection/<int:user_id>', methods=['POST'])
@login_required
@admin_required
def toggle_user_protection(user_id):
    """Activer/Désactiver la protection d'un utilisateur (réservé aux super admins)"""
    
    # Cette fonctionnalité pourrait être réservée à certains admins seulement
    if session.get('nom_compte') != 'skander.chniti':
        flash('Seul l\'administrateur principal peut modifier les protections', 'error')
        return redirect(url_for('admin_users'))
    
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    
    if not user:
        flash('Utilisateur introuvable', 'error')
        conn.close()
        return redirect(url_for('admin_users'))
    
    # Basculer le statut de protection
    new_status = 0 if user['is_protected'] == 1 else 1
    status_text = "activée" if new_status == 1 else "désactivée"
    
    conn.execute('''
        UPDATE users
        SET is_protected = ?, modifie_par = ?, updated_at = CURRENT_TIMESTAMP
        WHERE id = ?
    ''', (new_status, session['user_name'], user_id))
    
    conn.commit()
    conn.close()
    
    flash(f'Protection {status_text} pour {user["prenom"]} {user["nom"]}', 'success')
    return redirect(url_for('admin_users'))

@app.route('/admin/thresholds')
@login_required
@admin_required
def admin_thresholds():
    """Gestion des seuils"""
    conn = get_db_connection()
    seuils = conn.execute('''
        SELECT s1.* FROM seuils s1
        INNER JOIN (
            SELECT atelier, indicateur, MAX(date_modification) as max_date
            FROM seuils
            GROUP BY atelier, indicateur
        ) s2 ON s1.atelier = s2.atelier 
            AND s1.indicateur = s2.indicateur 
            AND s1.date_modification = s2.max_date
        ORDER BY s1.atelier, s1.indicateur
    ''').fetchall()
    conn.close()
    
    # Organiser les seuils par atelier
    seuils_by_atelier = {}
    for seuil in seuils:
        atelier = seuil['atelier']
        if atelier not in seuils_by_atelier:
            seuils_by_atelier[atelier] = {}
        
        indicateur = seuil['indicateur']
        seuils_by_atelier[atelier][indicateur] = seuil
    
    return render_template('admin/thresholds.html', seuils_by_atelier=seuils_by_atelier)

@app.route('/admin/update_threshold', methods=['POST'])
@login_required
@admin_required
def update_threshold():
    """Mettre à jour un seuil"""
    data = request.form
    
    conn = get_db_connection()
    
    # Vérifier si le seuil existe déjà
    existing_threshold = conn.execute('''
        SELECT id FROM seuils 
        WHERE atelier = ? AND indicateur = ?
    ''', (data['atelier'], data['indicateur'])).fetchone()
    
    if existing_threshold:
        # Mettre à jour la ligne existante
        conn.execute('''
            UPDATE seuils 
            SET valeur_seuil = ?, modifie_by = ?, date_modification = CURRENT_TIMESTAMP
            WHERE atelier = ? AND indicateur = ?
        ''', (float(data['valeur_seuil']), session['user_name'], 
              data['atelier'], data['indicateur']))
        
        flash('Seuil mis à jour avec succès', 'success')
    else:
        # Créer un nouveau seuil seulement s'il n'existe pas
        conn.execute('''
            INSERT INTO seuils (atelier, indicateur, valeur_seuil, modifie_by, date_modification)
            VALUES (?, ?, ?, ?, CURRENT_TIMESTAMP)
        ''', (data['atelier'], data['indicateur'], 
              float(data['valeur_seuil']), session['user_name']))
        
        flash('Nouveau seuil créé avec succès', 'success')
    
    conn.commit()
    conn.close()
    
    return redirect(url_for('admin_thresholds'))

# Routes pour le chargé qualité
@app.route('/admin/consultation')
@login_required
def quality_consultation():
    """Page de consultation des données (remplace data_view)"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        flash('Accès non autorisé', 'error')
        return redirect(url_for('dashboard'))
    
    return render_template('admin/consultation.html')


# API pour les statistiques du dashboard
@app.route('/api/dashboard_stats')
@login_required
def get_dashboard_stats():
    """Obtenir les statistiques du dashboard"""
    conn = get_db_connection()
    
    # Compter les utilisateurs actifs
    total_users = conn.execute('SELECT COUNT(*) as count FROM users').fetchone()['count']
    
    # Compter les seuils configurés
    total_thresholds = conn.execute('''
        SELECT COUNT(DISTINCT atelier || '-' || indicateur) as count FROM seuils
    ''').fetchone()['count']
    
    # Compter les rapports de ce mois
    current_month = datetime.now().month
    current_year = datetime.now().year
    
    monthly_reports = 0
    tables = ['pieces_retouchees', 'pieces_rebutees', 'pieces_controlees']
    for table in tables:
        count = conn.execute(f'''
            SELECT COUNT(*) as count FROM {table}
            WHERE mois = ? AND annee = ?
        ''', (current_month, current_year)).fetchone()['count']
        monthly_reports += count
    
    # Activité récente
    recent_activities = []
    
    # Derniers utilisateurs créés
    recent_users = conn.execute('''
        SELECT nom, prenom, created_at FROM users 
        ORDER BY created_at DESC LIMIT 3
    ''').fetchall()
    
    for user in recent_users:
        time_diff = datetime.now() - datetime.fromisoformat(user['created_at'])
        if time_diff.days == 0:
            time_str = "Aujourd'hui"
        elif time_diff.days == 1:
            time_str = "Hier"
        else:
            time_str = f"Il y a {time_diff.days} jours"
            
        recent_activities.append({
            'type': 'user',
            'icon': 'fas fa-user-plus',
            'text': f'Nouvel utilisateur: {user["prenom"]} {user["nom"]}',
            'time': time_str
        })
    
    # Derniers seuils modifiés
    recent_thresholds = conn.execute('''
        SELECT atelier, indicateur, modifie_by, date_modification FROM seuils 
        ORDER BY date_modification DESC LIMIT 2
    ''').fetchall()
    
    for threshold in recent_thresholds:
        time_diff = datetime.now() - datetime.fromisoformat(threshold['date_modification'])
        if time_diff.days == 0:
            time_str = "Aujourd'hui"
        elif time_diff.days == 1:
            time_str = "Hier"
        else:
            time_str = f"Il y a {time_diff.days} jours"
            
        recent_activities.append({
            'type': 'threshold',
            'icon': 'fas fa-sliders-h',
            'text': f'Seuil modifié: {threshold["indicateur"]} - {threshold["atelier"]}',
            'time': time_str
        })
    
    conn.close()
    
    return jsonify({
        'total_users': total_users,
        'total_thresholds': total_thresholds,
        'monthly_reports': monthly_reports,
        'recent_activities': recent_activities[:5]  # Limiter à 5 activités
    })

def calculate_ytd(conn, ateliers_interface, annee, mois_max, uap, atelier_param):
    total_pieces_exportees = total_pieces_reclamees = total_alertes = 0
    total_pieces_retouchees = total_pieces_rebutees = total_pieces_controlees = 0
    total_cout_rebut = total_reclamations_off = 0

    uap_condition = " AND uap = ?" if uap != 'all' else ""
    uap_params = [uap] if uap != 'all' else []

    for mois in range(1, mois_max + 1):
        for atelier in ateliers_interface:
            base_params = [mois, annee, atelier] + uap_params
            
            total_pieces_exportees += conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) FROM pieces_exportees 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_pieces_reclamees += conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) FROM pieces_reclamees 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_alertes += conn.execute(f'''
                SELECT COALESCE(SUM(nombre_alertes), 0) FROM quantite_alertes 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_pieces_retouchees += conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) FROM pieces_retouchees 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_pieces_rebutees += conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) FROM pieces_rebutees 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_pieces_controlees += conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) FROM pieces_controlees 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_cout_rebut += conn.execute(f'''
                SELECT COALESCE(SUM(montant), 0) FROM cout_rebut 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]
            
            total_reclamations_off += conn.execute(f'''
                SELECT COALESCE(SUM(nombre), 0) FROM reclamations_officielles 
                WHERE mois=? AND annee=? AND atelier=?{uap_condition}
            ''', base_params).fetchone()[0]

    ppm_officiel = (total_pieces_reclamees / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
    ppm_non_officiel = (total_alertes / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
    taux_rebut = (total_pieces_rebutees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0
    taux_retouche = (total_pieces_retouchees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0

    return {
        'ppm_officiel': ppm_officiel,
        'ppm_non_officiel': ppm_non_officiel,
        'taux_rebut': taux_rebut,
        'taux_retouche': taux_retouche,
        'cout_rebut': total_cout_rebut,
        'nombre_reclamations': total_reclamations_off,
        'cnq': 0
    }
# CODE CORRIGÉ POUR L'API CNQ

# Dans votre fichier app.py ou principal Flask, remplacez la route /api/charts_data par cette version corrigée et complète.
# Assurez-vous que les autres parties du code (connexion DB, etc.) restent inchangées.

# API Charts Data Corrigée - Section CNQ
# Modification à ajouter dans votre API get_charts_data

@app.route('/api/charts_data')
@login_required
def get_charts_data():
    atelier_param = request.args.get('atelier')
    annee = request.args.get('annee', type=int, default=datetime.now().year)
    uap = request.args.get('uap', 'all')
    aggregate = request.args.get('aggregate')
    comparison = request.args.get('comparison', 'false').lower() == 'true'
    raw_data = request.args.get('raw_data', 'false').lower() == 'true'  # NOUVEAU PARAMÈTRE
    
    # Si comparaison demandée, rediriger vers l'API de comparaison
    if comparison:
        return get_charts_data_comparison()

    conn = get_db_connection()
    
    # Liste des ateliers individuels
    all_ateliers = ['Colliers', 'Composite', 'Isolant souple', 'Manchons', 'Moulage', 'Protections thermiques', 'Racks', 'Système de visualisation']
    
    # Mapping des indicateurs pour les seuils
    seuil_mapping = {
        'ppm_officiel': 'PPM Officiel',
        'ppm_non_officiel': 'PPM Non Officiel',
        'taux_rebut': 'Taux de Rebut',
        'taux_retouche': 'Taux de Retouche',
        'cout_rebut': 'Coût de Rebut',
        'nombre_reclamations': 'Nombre de reclamation'
    }

    try:
        # Déterminer correctement le nom pour récupérer les seuils
        seuil_atelier_name = None
        
        # Déterminer les ateliers et le nom pour les seuils
        if not atelier_param or atelier_param == 'Total':
            ateliers_interface = all_ateliers
            seuil_atelier_name = 'Total'
            aggregate = 'sum'
        elif atelier_param == 'UAP1':
            ateliers_interface = ['Protections thermiques', 'Isolant souple']
            seuil_atelier_name = 'UAP1'
            aggregate = 'sum'
        elif atelier_param == 'UAP2':
            ateliers_interface = ['Manchons', 'Colliers', 'Racks', 'Moulage', 'Composite', 'Système de visualisation']
            seuil_atelier_name = 'UAP2'
            aggregate = 'sum'
        else:
            # Cas où atelier_param contient une liste d'ateliers séparés par des virgules
            ateliers_interface = [a.strip() for a in atelier_param.split(',')]
            
            # Détecter si c'est UAP1 ou UAP2 basé sur les ateliers
            uap1_ateliers = set(['Protections thermiques', 'Isolant souple'])
            uap2_ateliers = set(['Manchons', 'Colliers', 'Racks', 'Moulage', 'Composite', 'Système de visualisation'])
            ateliers_set = set(ateliers_interface)
            
            if ateliers_set == uap1_ateliers:
                seuil_atelier_name = 'UAP1'
                aggregate = 'sum'
            elif ateliers_set == uap2_ateliers:
                seuil_atelier_name = 'UAP2'
                aggregate = 'sum'
            elif len(ateliers_interface) == 1:
                seuil_atelier_name = ateliers_interface[0]
            else:
                seuil_atelier_name = None

        # Construction des conditions UAP
        uap_condition = ""
        uap_params = []
        if uap != 'all':
            uap_condition = " AND uap = ?"
            uap_params = [uap]

        monthly_data = {}
        raw_monthly_data = {}  # NOUVEAU : pour stocker les données brutes

        for mois in range(1, 13):
            # Initialiser les totaux
            total_pieces_exportees = 0
            total_pieces_reclamees = 0
            total_alertes = 0
            total_pieces_retouchees = 0
            total_pieces_rebutees = 0
            total_pieces_controlees = 0
            total_cout_rebut = 0
            total_reclamations_off = 0
            cnq_value = 0

            # Gestion des données CNQ depuis la table cnq (pour Total uniquement)
            if not atelier_param or atelier_param == 'Total' or aggregate == 'total':
                try:
                    cnq_result = conn.execute('''
                        SELECT valeur 
                        FROM cnq 
                        WHERE mois = ? AND annee = ? AND atelier = 'Total'
                    ''', (mois, annee)).fetchone()
                    
                    if cnq_result:
                        cnq_value = float(cnq_result['valeur'])
                    else:
                        cnq_value = 0
                    
                except Exception as e:
                    print(f"Erreur lors de la récupération des données CNQ pour mois {mois}: {e}")
                    cnq_value = 0
            else:
                cnq_value = 0

            # Pour les autres indicateurs, somme si aggregate ou plusieurs ateliers
            if aggregate == 'sum' or len(ateliers_interface) > 1:
                for atelier_interface in ateliers_interface:
                    base_params = [mois, annee, atelier_interface] + uap_params
                    
                    total_pieces_exportees += conn.execute(f'''
                        SELECT COALESCE(SUM(quantite), 0) as total 
                        FROM pieces_exportees
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_pieces_reclamees += conn.execute(f'''
                        SELECT COALESCE(SUM(quantite), 0) as total 
                        FROM pieces_reclamees
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_alertes += conn.execute(f'''
                        SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                        FROM quantite_alertes
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_pieces_retouchees += conn.execute(f'''
                        SELECT COALESCE(SUM(quantite), 0) as total 
                        FROM pieces_retouchees 
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_pieces_rebutees += conn.execute(f'''
                        SELECT COALESCE(SUM(quantite), 0) as total 
                        FROM pieces_rebutees 
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_pieces_controlees += conn.execute(f'''
                        SELECT COALESCE(SUM(quantite), 0) as total 
                        FROM pieces_controlees 
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_cout_rebut += conn.execute(f'''
                        SELECT COALESCE(SUM(montant), 0) as total 
                        FROM cout_rebut 
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
                    
                    total_reclamations_off += conn.execute(f'''
                        SELECT COALESCE(SUM(nombre), 0) as total 
                        FROM reclamations_officielles 
                        WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                    ''', base_params).fetchone()['total']
            else:
                # Pour un atelier unique
                atelier_interface = ateliers_interface[0]
                base_params = [mois, annee, atelier_interface] + uap_params
                
                total_pieces_exportees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_exportees 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_pieces_reclamees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_reclamees 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_alertes = conn.execute(f'''
                    SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                    FROM quantite_alertes 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_pieces_retouchees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_retouchees 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_pieces_rebutees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_rebutees 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_pieces_controlees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_controlees 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_cout_rebut = conn.execute(f'''
                    SELECT COALESCE(SUM(montant), 0) as total 
                    FROM cout_rebut 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']
                
                total_reclamations_off = conn.execute(f'''
                    SELECT COALESCE(SUM(nombre), 0) as total 
                    FROM reclamations_officielles 
                    WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                ''', base_params).fetchone()['total']

            # NOUVEAU : Stocker les données brutes si demandées
            if raw_data:
                raw_monthly_data[mois] = {
                    'pieces_exportees': total_pieces_exportees,
                    'pieces_reclamees': total_pieces_reclamees,
                    'alertes': total_alertes,
                    'pieces_retouchees': total_pieces_retouchees,
                    'pieces_rebutees': total_pieces_rebutees,
                    'pieces_controlees': total_pieces_controlees,
                    'cout_rebut': total_cout_rebut,
                    'reclamations_officielles': total_reclamations_off
                }

            # Calculer les indicateurs
            ppm_officiel = (total_pieces_reclamees / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
            ppm_non_officiel = (total_alertes / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
            taux_retouche = (total_pieces_retouchees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0
            taux_rebut = (total_pieces_rebutees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0

            monthly_data[mois] = {
                'ppm_officiel': ppm_officiel,
                'ppm_non_officiel': ppm_non_officiel,
                'cout_rebut':total_cout_rebut,
                'nombre_reclamations': total_reclamations_off,
                'cnq': cnq_value, 
                'taux_rebut': taux_rebut, 
                'taux_retouche': taux_retouche
            }

        # Récupérer les seuils correctement
        seuils = {}
        
        if seuil_atelier_name:
            # Pour tous les indicateurs standards depuis la table seuils
            for api_indicator, db_indicator in seuil_mapping.items():
                try:
                    seuil_result = conn.execute('''
                        SELECT valeur_seuil
                        FROM seuils
                        WHERE atelier = ? AND indicateur = ?
                        ORDER BY date_modification DESC
                        LIMIT 1
                    ''', (seuil_atelier_name, db_indicator)).fetchone()
                    
                    if seuil_result:
                        seuils[api_indicator] = float(seuil_result['valeur_seuil'])
                    else:
                        seuils[api_indicator] = 0
                        
                except Exception as e:
                    print(f"Erreur seuil {api_indicator} pour {seuil_atelier_name}: {e}")
                    seuils[api_indicator] = 0
        else:
            # Si pas d'atelier spécifique, mettre tous les seuils à 0
            for api_indicator in seuil_mapping.keys():
                seuils[api_indicator] = 0

        # Seuil pour CNQ depuis la table cnq (pour Total uniquement)
        if not atelier_param or atelier_param == 'Total' or aggregate == 'total':
            try:
                seuil_cnq = conn.execute('''
                    SELECT seuil 
                    FROM cnq 
                    WHERE atelier = 'Total' 
                    ORDER BY created_at DESC 
                    LIMIT 1
                ''').fetchone()
                if seuil_cnq:
                    seuils['cnq'] = float(seuil_cnq['seuil'])
                else:
                    seuils['cnq'] = 15.0  # Valeur par défaut
            except Exception as e:
                print(f"Erreur seuil CNQ: {e}")
                seuils['cnq'] = 15.0  # Valeur par défaut
        else:
            seuils['cnq'] = 0

        # Préparer la réponse
        response_data = {
            'monthly_data': monthly_data,
            'seuils': seuils,
            'annee': annee,
            'atelier': atelier_param or 'Total',
            'uap': uap,
            'aggregate': aggregate
        }
        
        # NOUVEAU : Ajouter les données brutes si demandées
        if raw_data:
            response_data['raw_monthly_data'] = raw_monthly_data
            current_year = datetime.now().year
            current_month = datetime.now().month - 1
            ytd_current = calculate_ytd(conn, ateliers_interface, current_year, current_month, uap, atelier_param)
            ytd_previous = calculate_ytd(conn, ateliers_interface, current_year - 1, current_month, uap, atelier_param)
            total_previous = calculate_ytd(conn, ateliers_interface, current_year - 1, 12, uap, atelier_param)
            response_data.update({
        'ytd_current': ytd_current,
        'ytd_previous': ytd_previous,
        'total_previous': total_previous
    })
        conn.close()
        return jsonify(response_data)

    except Exception as e:
        conn.close()
        print(f"Erreur générale dans get_charts_data: {e}")
        return jsonify({'error': f'Erreur serveur: {str(e)}'}), 500

# API YTD Data Corrigée
@app.route('/api/ytd_data')
@login_required
def get_ytd_data():
    atelier_param = request.args.get('atelier')
    annee = request.args.get('annee', type=int, default=datetime.now().year)
    uap = request.args.get('uap', 'all')
    
    conn = get_db_connection()
    
    # Liste des ateliers
    all_ateliers = ['Colliers', 'Composite', 'Isolant souple', 'Manchons', 'Moulage', 'Protections thermiques', 'Racks', 'Système de visualisation']
    
    try:
        # Déterminer les ateliers à traiter
        if atelier_param in ['UAP1', 'UAP2', 'Total']:
            if atelier_param == 'UAP1':
                ateliers_interface = ['Protections thermiques', 'Isolant souple']
            elif atelier_param == 'UAP2':
                ateliers_interface = ['Manchons', 'Colliers', 'Racks', 'Moulage', 'Composite', 'Système de visualisation']
            else:  # Total
                ateliers_interface = all_ateliers
        else:
            ateliers_interface = all_ateliers
        
        # UAP condition
        uap_condition = " AND uap = ?" if uap != 'all' else ""
        uap_params = [uap] if uap != 'all' else []
        
        # Période YTD (MOIS ACTUEL - 1)
        if annee == datetime.now().year:
            current_month = max(datetime.now().month - 1, 1)  # OCTOBRE → 9
        else:
            current_month = 12
        
        # Initialiser totaux
        ytd_pieces_exportees = ytd_pieces_reclamees = ytd_alertes = 0
        ytd_pieces_retouchees = ytd_pieces_rebutees = ytd_pieces_controlees = 0
        ytd_cout_rebut = ytd_reclamations_off = 0
        ytd_cnq_sum = 0  # ✅ SOMME CNQ
        
        # BOUCLE MOIS
        for mois in range(1, current_month + 1):
            # ✅ CNQ : UNIQUEMENT POUR atelier_param == 'Total'
            if atelier_param == 'Total':
                cnq_result = conn.execute('''
                    SELECT valeur 
                    FROM cnq 
                    WHERE mois = ? AND annee = ? AND atelier = 'Total'
                ''', (mois, annee)).fetchone()
                
                if cnq_result:
                    ytd_cnq_sum += cnq_result['valeur']
            
            # AUTRES INDICATEURS (tous ateliers)
            for atelier in ateliers_interface:
                base_params = [mois, annee, atelier] + uap_params
                
                ytd_pieces_exportees += conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) FROM pieces_exportees 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_pieces_reclamees += conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) FROM pieces_reclamees 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_alertes += conn.execute(f'''
                    SELECT COALESCE(SUM(nombre_alertes), 0) FROM quantite_alertes 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_pieces_retouchees += conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) FROM pieces_retouchees 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_pieces_rebutees += conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) FROM pieces_rebutees 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_pieces_controlees += conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) FROM pieces_controlees 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_cout_rebut += conn.execute(f'''
                    SELECT COALESCE(SUM(montant), 0) FROM cout_rebut 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
                
                ytd_reclamations_off += conn.execute(f'''
                    SELECT COALESCE(SUM(nombre), 0) FROM reclamations_officielles 
                    WHERE mois=? AND annee=? AND atelier=?{uap_condition}
                ''', base_params).fetchone()[0]
        
        # CALCULS FINAUX
        ytd_ppm_officiel = (ytd_pieces_reclamees / ytd_pieces_exportees * 1000000) if ytd_pieces_exportees > 0 else 0
        ytd_ppm_non_officiel = (ytd_alertes / ytd_pieces_exportees * 1000000) if ytd_pieces_exportees > 0 else 0
        ytd_taux_retouche = (ytd_pieces_retouchees / ytd_pieces_controlees * 100) if ytd_pieces_controlees > 0 else 0
        ytd_taux_rebut = (ytd_pieces_rebutees / ytd_pieces_controlees * 100) if ytd_pieces_controlees > 0 else 0
        ytd_cnq = ytd_cnq_sum  # ✅ SOMME PURE = 4.89
        
        display_name = atelier_param or 'Total'
        
        ytd_data = {
            display_name: {
                'ppm_officiel': round(ytd_ppm_officiel, 3),
                'ppm_non_officiel': round(ytd_ppm_non_officiel, 3),
                'taux_rebut': round(ytd_taux_rebut, 3),
                'taux_retouche': round(ytd_taux_retouche, 3),
                'cout_rebut': round(ytd_cout_rebut, 3),
                'cnq': round(ytd_cnq, 3),  # ✅ 4.890
                'nombre_reclamations': int(ytd_reclamations_off)
            }
        }
        
        conn.close()
        
        # DEBUG
        print(f"=== YTD {display_name} {annee} ===")
        print(f"CNQ SOMME: {ytd_cnq} (9 mois)")
        print(f"PPM: {ytd_ppm_officiel}")
        print("========================")
        
        return jsonify({'ytd_data': ytd_data, 'annee': annee, 'uap': uap})
        
    except Exception as e:
        conn.close()
        return jsonify({'error': str(e)}), 500



# Nouvelle API pour obtenir les données de conformité
@app.route('/api/conformity_data')
@login_required  
def get_conformity_data():
    atelier = request.args.get('atelier')
    annee = request.args.get('annee', type=int, default=2025)
    
    if not atelier:
        return jsonify({'error': 'Atelier requis'}), 400
    
    conn = get_db_connection()
    
    # Mapping des noms d'ateliers
    atelier_mapping = {
        'Manchon': ['Manchons'],
        'Collier de Fixation': ['Colliers'],
        'Rack': ['Racks'],
        'Moulage': ['Moulage'],
        'Isolation Thermique': ['Protections thermiques', 'Protection thermique', 'Isolant thermique'],
        'Isolation Souple': ['Isolant souple'],
        'Composite': ['Composite'],
        'Atelier de Visualisation': ['Système de visualisation']
    }
    
    # Mapping des indicateurs pour les seuils
    seuil_mapping = {
        'ppm_officiel': 'PPM Officiel',
        'ppm_non_officiel': 'PPM Non Officiel', 
        'taux_rebut': 'Taux de Rebut',
        'taux_retouche': 'Taux de Retouche',
        'cout_rebut': 'Coût de Rebut',
        'cnq': 'CNQ',
        'nombre_reclamations': 'Nombre de reclamation'
    }
    
    try:
        # Gestion des cas spéciaux
        if atelier == 'UAP 1':
            ateliers_interface = ['Isolation Thermique', 'Isolation Souple']
        elif atelier == 'UAP 2':
            ateliers_interface = ['Manchon', 'Collier de Fixation', 'Rack', 'Moulage', 'Composite', 'Atelier de Visualisation']
        elif atelier == 'Total':
            ateliers_interface = list(atelier_mapping.keys())
        else:
            ateliers_interface = [atelier]
        
        # Convertir vers les noms de base de données
        ateliers_db = []
        for atelier_int in ateliers_interface:
            if atelier_int in atelier_mapping:
                ateliers_db.extend(atelier_mapping[atelier_int])
            else:
                ateliers_db.append(atelier_int)
        
        atelier_condition = f"atelier IN ({','.join(['?' for _ in ateliers_db])})"
        
        # Obtenir les données mensuelles
        monthly_data = {}
        current_month = 8  # Août 2025
        
        for mois in range(1, current_month + 1):
            month_params = [mois, annee] + ateliers_db
            
            try:
                # Récupérer les données comme dans charts_data
                pieces_exportees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_exportees 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                pieces_reclamees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_reclamees 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                alertes = conn.execute(f'''
                    SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                    FROM quantite_alertes 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                pieces_retouchees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_retouchees 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                pieces_rebutees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total 
                    FROM pieces_rebutees 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                pieces_controlees = conn.execute(f'''
                    SELECT COALESCE(SUM(quantite), 0) as total  
                    FROM pieces_controlees 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                cout_rebut = conn.execute(f'''
                    SELECT COALESCE(SUM(montant), 0) as total 
                    FROM cout_rebut 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                cnq_val = conn.execute(f'''
                    SELECT COALESCE(AVG(valeur), 0) as avg 
                    FROM cnq 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['avg']

                reclamations_off = conn.execute(f'''
                    SELECT COALESCE(SUM(nombre), 0) as total 
                    FROM reclamations_officielles 
                    WHERE mois = ? AND annee = ? AND {atelier_condition}
                ''', month_params).fetchone()['total']

                # Calculer les indicateurs
                ppm_officiel = (pieces_reclamees / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
                ppm_non_officiel = (alertes / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
                taux_retouche = (pieces_retouchees / pieces_controlees * 100) if pieces_controlees > 0 else 0
                taux_rebut = (pieces_rebutees / pieces_controlees * 100) if pieces_controlees > 0 else 0

                monthly_data[mois] = {
                    'ppm_officiel': round(ppm_officiel, 2),
                    'ppm_non_officiel': round(ppm_non_officiel, 2),
                    'cout_rebut': round(cout_rebut, 2),
                    'nombre_reclamations': reclamations_off,
                    'cnq': round(cnq_val, 2),
                    'taux_rebut': round(taux_rebut, 2),
                    'taux_retouche': round(taux_retouche, 2)
                }

            except Exception as e:
                print(f"Erreur pour le mois {mois}: {e}")
                monthly_data[mois] = {
                    'ppm_officiel': 0,
                    'ppm_non_officiel': 0,
                    'cout_rebut': 0,
                    'nombre_reclamations': 0,
                    'cnq': 0,
                    'taux_rebut': 0,
                    'taux_retouche': 0
                }
        
        # Récupérer les seuils
        seuils = {}
        
        if len(ateliers_interface) > 1:
            # Pour agrégats: moyenne des seuils
            for api_indicator, db_indicator in seuil_mapping.items():
                try:
                    seuils_placeholders = ','.join(['?' for _ in ateliers_interface])
                    seuil_result = conn.execute(f'''
                        SELECT AVG(valeur_seuil) as valeur_seuil
                        FROM seuils 
                        WHERE atelier IN ({seuils_placeholders}) AND indicateur = ?
                    ''', ateliers_interface + [db_indicator]).fetchone()
                    
                    if seuil_result and seuil_result['valeur_seuil']:
                        seuils[api_indicator] = float(seuil_result['valeur_seuil'])
                    else:
                        seuils[api_indicator] = 0
                except Exception as e:
                    print(f"Erreur seuil {api_indicator}: {e}")
                    seuils[api_indicator] = 0
        else:
            # Pour un seul atelier
            for api_indicator, db_indicator in seuil_mapping.items():
                try:
                    seuil_result = conn.execute('''
                        SELECT valeur_seuil 
                        FROM seuils
                        WHERE atelier = ? AND indicateur = ?
                        ORDER BY date_modification DESC
                        LIMIT 1
                    ''', (atelier, db_indicator)).fetchone()
                    
                    if seuil_result:
                        seuils[api_indicator] = float(seuil_result['valeur_seuil'])
                    else:
                        seuils[api_indicator] = 0
                except Exception as e:
                    print(f"Erreur seuil {api_indicator}: {e}")
                    seuils[api_indicator] = 0
        
        # Analyser la conformité
        conformity_analysis = {}
        indicators = ['ppm_officiel', 'ppm_non_officiel', 'taux_rebut', 'taux_retouche', 'cout_rebut', 'cnq', 'nombre_reclamations']
        
        for indicator in indicators:
            threshold = seuils.get(indicator, 0)
            conformity_analysis[indicator] = {
                'threshold': threshold,
                'months': {}
            }
            
            for mois in range(1, current_month + 1):
                value = monthly_data[mois][indicator]
                is_conform = value <= threshold if threshold > 0 else True
                conformity_analysis[indicator]['months'][mois] = {
                    'value': value,
                    'is_conform': is_conform,
                    'status': 'Conforme' if is_conform else 'Non Conforme'
                }
        
        conn.close()
        
        return jsonify({
            'monthly_data': monthly_data,
            'seuils': seuils,
            'conformity_analysis': conformity_analysis,
            'atelier': atelier,
            'current_month': current_month
        })
        
    except Exception as e:
        conn.close()
        print(f"Erreur dans get_conformity_data: {e}")
        return jsonify({'error': f'Erreur serveur: {str(e)}'}), 500
# Gestion d'erreurs
@app.errorhandler(404)
def not_found_error(error):
    return render_template('errors/404.html'), 404

@app.errorhandler(500)
def internal_error(error):
    return render_template('errors/500.html'), 500

@app.errorhandler(403)
def forbidden_error(error):
    return render_template('errors/403.html'), 403

# Route de test pour vérifier la base de données
@app.route('/api/test_db')
@login_required
def test_db():
    """Test de la base de données"""
    if session.get('user_role') != 'Administration':
        return jsonify({'error': 'Accès non autorisé'}), 403
    
    conn = get_db_connection()
    
    # Compter les enregistrements dans chaque table
    tables = [
        'users', 'seuils', 'pieces_retouchees', 'pieces_rebutees',
        'pieces_controlees', 'cout_rebut', 'cnq', 'pieces_reclamees',
        'pieces_exportees', 'quantite_alertes', 'reclamations_officielles'
    ]
    
    counts = {}
    for table in tables:
        try:
            count = conn.execute(f'SELECT COUNT(*) as count FROM {table}').fetchone()['count']
            counts[table] = count
        except Exception as e:
            counts[table] = f'Erreur: {str(e)}'
    
    conn.close()
    
    return jsonify({
        'database_status': 'OK',
        'table_counts': counts,
        'timestamp': datetime.now().isoformat()
    })



@app.route('/quality/submit_data', methods=['POST'])
@login_required
def submit_quality_data():
    """Traiter la soumission des données qualité"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        flash('Accès non autorisé', 'error')
        return redirect(url_for('dashboard'))
    
    try:
        # Récupérer les données du formulaire
        data = request.form
        
        mois = int(data['mois'])
        annee = int(data['annee'])
        atelier = data['atelier']
        uap = data['uap']
        
        # Données qualité interne
        pieces_retouchees = int(data['pieces_retouchees'])
        pieces_rebutees = int(data['pieces_rebutees'])
        pieces_controlees = int(data['pieces_controlees'])
        cout_rebut = float(data['cout_rebut'])
        cnq = float(data['cnq'])
        
        # Données qualité externe
        pieces_reclamees = int(data['pieces_reclamees'])
        pieces_exportees = int(data['pieces_exportees'])
        quantite_alertes = int(data['quantite_alertes'])
        reclamations_officielles = int(data['reclamations_officielles'])
        
        conn = get_db_connection()
        
        # Insérer ou mettre à jour chaque type de données
        # Pièces retouchées
        conn.execute('''
            INSERT OR REPLACE INTO pieces_retouchees 
            (mois, annee, atelier, uap, quantite)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, pieces_retouchees))
        
        # Pièces rebutées
        conn.execute('''
            INSERT OR REPLACE INTO pieces_rebutees 
            (mois, annee, atelier, uap, quantite)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, pieces_rebutees))
        
        # Pièces contrôlées
        conn.execute('''
            INSERT OR REPLACE INTO pieces_controlees 
            (mois, annee, atelier, uap, quantite)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, pieces_controlees))
        
        # Coût du rebut
        conn.execute('''
            INSERT OR REPLACE INTO cout_rebut 
            (mois, annee, atelier, uap, montant)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, cout_rebut))
        
        # CNQ
        conn.execute('''
            INSERT OR REPLACE INTO cnq 
            (mois, annee, atelier, uap, valeur)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, cnq))
        
        # Pièces réclamées
        conn.execute('''
            INSERT OR REPLACE INTO pieces_reclamees 
            (mois, annee, atelier, uap, quantite)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, pieces_reclamees))
        
        # Pièces exportées
        conn.execute('''
            INSERT OR REPLACE INTO pieces_exportees 
            (mois, annee, atelier, uap, quantite)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, pieces_exportees))
        
        # Quantité d'alertes
        conn.execute('''
            INSERT OR REPLACE INTO quantite_alertes 
            (mois, annee, atelier, uap, nombre_alertes)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, quantite_alertes))
        
        # Réclamations officielles
        conn.execute('''
            INSERT OR REPLACE INTO reclamations_officielles 
            (mois, annee, atelier, uap, nombre)
            VALUES (?, ?, ?, ?, ?)
        ''', (mois, annee, atelier, uap, reclamations_officielles))
        
        conn.commit()
        conn.close()
        
        # Calculer les indicateurs pour affichage
        ppm_officiel = (pieces_reclamees / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
        ppm_non_officiel = (quantite_alertes / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
        taux_retouche = (pieces_retouchees / pieces_controlees * 100) if pieces_controlees > 0 else 0
        taux_rebut = (pieces_rebutees / pieces_controlees * 100) if pieces_controlees > 0 else 0
        
        flash(f'''Données enregistrées avec succès ! 
               PPM Officiel: {ppm_officiel:.2f}, 
               PPM Non Officiel: {ppm_non_officiel:.2f}, 
               Taux de Rebut: {taux_rebut:.2f}%, 
               Taux de Retouche: {taux_retouche:.2f}%''', 'success')
        
        return redirect(url_for('quality_data_entry'))
        
    except Exception as e:
        flash(f'Erreur lors de l\'enregistrement des données: {str(e)}', 'error')
        return redirect(url_for('quality_data_entry'))

# Route pour obtenir les données existantes (pour pré-remplir le formulaire)
@app.route('/api/get_existing_data')
@login_required
def get_existing_data():
    """Obtenir les données existantes pour pré-remplir le formulaire"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        return jsonify({'error': 'Accès non autorisé'}), 403
    
    mois = request.args.get('mois', type=int)
    annee = request.args.get('annee', type=int)
    atelier = request.args.get('atelier')
    uap = request.args.get('uap')
    
    if not all([mois, annee, atelier, uap]):
        return jsonify({'error': 'Paramètres manquants'}), 400
    
    conn = get_db_connection()
    result = {}
    
    # Récupérer toutes les données existantes
    tables_queries = {
        'pieces_retouchees': 'SELECT quantite FROM pieces_retouchees WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'pieces_rebutees': 'SELECT quantite FROM pieces_rebutees WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'pieces_controlees': 'SELECT quantite FROM pieces_controlees WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'cout_rebut': 'SELECT montant FROM cout_rebut WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'cnq': 'SELECT valeur FROM cnq WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'pieces_reclamees': 'SELECT quantite FROM pieces_reclamees WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'pieces_exportees': 'SELECT quantite FROM pieces_exportees WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'quantite_alertes': 'SELECT nombre_alertes FROM quantite_alertes WHERE mois=? AND annee=? AND atelier=? AND uap=?',
        'reclamations_officielles': 'SELECT nombre FROM reclamations_officielles WHERE mois=? AND annee=? AND atelier=? AND uap=?'
    }
    
    for table_name, query in tables_queries.items():
        data = conn.execute(query, (mois, annee, atelier, uap)).fetchone()
        if data:
            if table_name == 'cout_rebut':
                result[table_name] = data['montant']
            elif table_name == 'cnq':
                result[table_name] = data['valeur']
            elif table_name == 'quantite_alertes':
                result[table_name] = data['nombre_alertes']
            elif table_name == 'reclamations_officielles':
                result[table_name] = data['nombre']
            else:
                result[table_name] = data['quantite']
        else:
            result[table_name] = 0
    
    conn.close()
    
    return jsonify({
        'exists': any(result.values()),
        'data': result
    })

# Routes pour la consultation des données
@app.route('/api/consultation_data')
@login_required
def get_consultation_data():
    """Obtenir les données pour la consultation selon les critères"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        return jsonify({'error': 'Accès non autorisé'}), 403
    
    atelier = request.args.get('atelier')
    mois = request.args.get('mois', type=int)
    annee = request.args.get('annee', type=int)
    
    if not all([atelier, mois, annee]):
        return jsonify({'error': 'Tous les paramètres sont requis'}), 400
    
    conn = get_db_connection()
    
    # Déterminer les ateliers à inclure selon la sélection
    if atelier == 'UAP1':
        ateliers_list = ['Isolation Thermique', 'Isolation Souple']
    elif atelier == 'UAP2':
        ateliers_list = ['Manchon', 'Collier de Fixation', 'Rack', 'Moulage', 
                        'Composite', 'Atelier de Visualisation']
    elif atelier == 'Usine Complète':
        ateliers_list = ['Manchon', 'Collier de Fixation', 'Rack', 'Moulage',
                        'Isolation Thermique', 'Isolation Souple', 'Composite',
                        'Atelier de Visualisation']
    else:
        ateliers_list = [atelier]
    
    # Construire la condition WHERE pour les ateliers
    placeholders = ','.join(['?' for _ in ateliers_list])
    base_params = [mois, annee] + ateliers_list
    
    # Récupérer les données agrégées
    pieces_exportees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_exportees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_reclamees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_reclamees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    alertes = conn.execute(f'''
        SELECT COALESCE(SUM(nombre_alertes), 0) as total 
        FROM quantite_alertes 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_retouchees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_retouchees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_rebutees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_rebutees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_controlees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_controlees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    cout_rebut = conn.execute(f'''
        SELECT COALESCE(SUM(montant), 0) as total 
        FROM cout_rebut 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    cnq_val = conn.execute(f'''
        SELECT COALESCE(AVG(valeur), 0) as avg 
        FROM cnq 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['avg']
    
    reclamations_off = conn.execute(f'''
        SELECT COALESCE(SUM(nombre), 0) as total 
        FROM reclamations_officielles 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    # Calculer les indicateurs
    ppm_officiel = (pieces_reclamees / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
    ppm_non_officiel = (alertes / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
    taux_retouche = (pieces_retouchees / pieces_controlees * 100) if pieces_controlees > 0 else 0
    taux_rebut = (pieces_rebutees / pieces_controlees * 100) if pieces_controlees > 0 else 0
    
    # Récupérer les seuils (prendre la moyenne des seuils des ateliers concernés)
    seuils = {}
    if len(ateliers_list) == 1:
        # Un seul atelier, prendre ses seuils
        seuils_data = conn.execute('''
            SELECT s1.indicateur, s1.valeur_seuil 
            FROM seuils s1
            INNER JOIN (
                SELECT atelier, indicateur, MAX(date_modification) as max_date
                FROM seuils
                WHERE atelier = ?
                GROUP BY atelier, indicateur
            ) s2 ON s1.atelier = s2.atelier 
                AND s1.indicateur = s2.indicateur 
                AND s1.date_modification = s2.max_date
        ''', (ateliers_list[0],)).fetchall()
    else:
        # Plusieurs ateliers, calculer la moyenne des seuils
        placeholders_seuils = ','.join(['?' for _ in ateliers_list])
        seuils_data = conn.execute(f'''
            SELECT s1.indicateur, AVG(s1.valeur_seuil) as valeur_seuil
            FROM seuils s1
            INNER JOIN (
                SELECT atelier, indicateur, MAX(date_modification) as max_date
                FROM seuils
                WHERE atelier IN ({placeholders_seuils})
                GROUP BY atelier, indicateur
            ) s2 ON s1.atelier = s2.atelier 
                AND s1.indicateur = s2.indicateur 
                AND s1.date_modification = s2.max_date
            GROUP BY s1.indicateur
        ''', ateliers_list).fetchall()
    
    for seuil in seuils_data:
        seuils[seuil['indicateur']] = seuil['valeur_seuil']
    
    # Récupérer les données détaillées par atelier si nécessaire
    details_ateliers = []
    if len(ateliers_list) > 1:
        for atelier_detail in ateliers_list:
            detail_params = [mois, annee, atelier_detail]
            
            pieces_exp_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_exportees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_recl_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_reclamees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            alertes_detail = conn.execute('''
                SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                FROM quantite_alertes 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_ret_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_retouchees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_reb_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_rebutees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_ctrl_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_controlees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            cout_reb_detail = conn.execute('''
                SELECT COALESCE(SUM(montant), 0) as total 
                FROM cout_rebut 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            cnq_detail = conn.execute('''
                SELECT COALESCE(AVG(valeur), 0) as avg 
                FROM cnq 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['avg']
            
            recl_off_detail = conn.execute('''
                SELECT COALESCE(SUM(nombre), 0) as total 
                FROM reclamations_officielles 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            # Calculer les indicateurs pour cet atelier
            ppm_off_detail = (pieces_recl_detail / pieces_exp_detail * 1000000) if pieces_exp_detail > 0 else 0
            ppm_non_off_detail = (alertes_detail / pieces_exp_detail * 1000000) if pieces_exp_detail > 0 else 0
            taux_ret_detail = (pieces_ret_detail / pieces_ctrl_detail * 100) if pieces_ctrl_detail > 0 else 0
            taux_reb_detail = (pieces_reb_detail / pieces_ctrl_detail * 100) if pieces_ctrl_detail > 0 else 0
            
            details_ateliers.append({
                'atelier': atelier_detail,
                'ppm_officiel': round(ppm_off_detail, 2),
                'ppm_non_officiel': round(ppm_non_off_detail, 2),
                'taux_retouche': round(taux_ret_detail, 2),
                'taux_rebut': round(taux_reb_detail, 2),
                'cout_rebut': round(cout_reb_detail, 2),
                'cnq': round(cnq_detail, 2),
                'nombre_reclamations': recl_off_detail,
                'pieces_exportees': pieces_exp_detail,
                'pieces_controlees': pieces_ctrl_detail
            })
    
    conn.close()
    
    # Préparer la réponse
    result = {
        'atelier_selection': atelier,
        'mois': mois,
        'annee': annee,
        'ateliers_inclus': ateliers_list,
        'donnees_globales': {
            'ppm_officiel': round(ppm_officiel, 2),
            'ppm_non_officiel': round(ppm_non_officiel, 2),
            'taux_retouche': round(taux_retouche, 2),
            'taux_rebut': round(taux_rebut, 2),
            'cout_rebut': round(cout_rebut, 2),
            'cnq': round(cnq_val, 2),
            'nombre_reclamations': reclamations_off,
            'pieces_exportees': pieces_exportees,
            'pieces_controlees': pieces_controlees,
            'pieces_reclamees': pieces_reclamees,
            'pieces_retouchees': pieces_retouchees,
            'pieces_rebutees': pieces_rebutees,
            'quantite_alertes': alertes
        },
        'seuils': seuils,
        'details_ateliers': details_ateliers
    }
    
    return jsonify(result)

@app.route('/api/export_consultation_excel')
@login_required  
def export_consultation_excel():
    """Exporter les données de consultation en Excel"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        return jsonify({'error': 'Accès non autorisé'}), 403
    
    atelier = request.args.get('atelier')
    mois = request.args.get('mois', type=int)
    annee = request.args.get('annee', type=int)
    
    if not all([atelier, mois, annee]):
        return jsonify({'error': 'Tous les paramètres sont requis'}), 400
    
    # Récupérer les données (réutiliser la logique de get_consultation_data)
    conn = get_db_connection()
    
    # Déterminer les ateliers à inclure
    if atelier == 'UAP1':
        ateliers_list = ['Isolation Thermique', 'Isolation Souple']
    elif atelier == 'UAP2':
        ateliers_list = ['Manchon', 'Collier de Fixation', 'Rack', 'Moulage', 
                        'Composite', 'Atelier de Visualisation']
    elif atelier == 'Usine Complète':
        ateliers_list = ['Manchon', 'Collier de Fixation', 'Rack', 'Moulage',
                        'Isolation Thermique', 'Isolation Souple', 'Composite',
                        'Atelier de Visualisation']
    else:
        ateliers_list = [atelier]
    
    # Créer un fichier CSV en mémoire (simulant Excel)
    output = io.StringIO()
    writer = csv.writer(output)
    
    # Noms des mois
    month_names = ['', 'Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
                   'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre']
    
    # En-tête du rapport
    writer.writerow([f'Rapport Qualité - {atelier}'])
    writer.writerow([f'Période: {month_names[mois]} {annee}'])
    writer.writerow([f'Ateliers inclus: {", ".join(ateliers_list)}'])
    writer.writerow([f'Date d\'export: {datetime.now().strftime("%d/%m/%Y %H:%M")}'])
    writer.writerow([])
    
    # Headers des données globales
    writer.writerow([
        'Type de Données', 'Valeur', 'Unité'
    ])
    
    # Récupérer et écrire les données globales
    placeholders = ','.join(['?' for _ in ateliers_list])
    base_params = [mois, annee] + ateliers_list
    
    # Calculer tous les indicateurs (même logique qu'avant)
    pieces_exportees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_exportees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_reclamees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_reclamees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    alertes = conn.execute(f'''
        SELECT COALESCE(SUM(nombre_alertes), 0) as total 
        FROM quantite_alertes 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_retouchees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_retouchees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_rebutees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_rebutees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    pieces_controlees = conn.execute(f'''
        SELECT COALESCE(SUM(quantite), 0) as total 
        FROM pieces_controlees 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    cout_rebut = conn.execute(f'''
        SELECT COALESCE(SUM(montant), 0) as total 
        FROM cout_rebut 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    cnq_val = conn.execute(f'''
        SELECT COALESCE(AVG(valeur), 0) as avg 
        FROM cnq 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['avg']
    
    reclamations_off = conn.execute(f'''
        SELECT COALESCE(SUM(nombre), 0) as total 
        FROM reclamations_officielles 
        WHERE mois = ? AND annee = ? AND atelier IN ({placeholders})
    ''', base_params).fetchone()['total']
    
    # Calculer les indicateurs
    ppm_officiel = (pieces_reclamees / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
    ppm_non_officiel = (alertes / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
    taux_retouche = (pieces_retouchees / pieces_controlees * 100) if pieces_controlees > 0 else 0
    taux_rebut = (pieces_rebutees / pieces_controlees * 100) if pieces_controlees > 0 else 0
    
    # Écrire les données
    donnees_export = [
        ('Pièces Exportées', pieces_exportees, 'pièces'),
        ('Pièces Contrôlées', pieces_controlees, 'pièces'),
        ('Pièces Réclamées', pieces_reclamees, 'pièces'),
        ('Pièces Retouchées', pieces_retouchees, 'pièces'),
        ('Pièces Rebutées', pieces_rebutees, 'pièces'),
        ('Quantité d\'Alertes', alertes, 'alertes'),
        ('Réclamations Officielles', reclamations_off, 'réclamations'),
        ('PPM Officiel', round(ppm_officiel, 2), 'PPM'),
        ('PPM Non Officiel', round(ppm_non_officiel, 2), 'PPM'),
        ('Taux de Retouche', round(taux_retouche, 2), '%'),
        ('Taux de Rebut', round(taux_rebut, 2), '%'),
        ('Coût de Rebut', round(cout_rebut, 2), 'TND'),
        ('CNQ', round(cnq_val, 2), 'TND')
    ]
    
    for donnee in donnees_export:
        writer.writerow(donnee)
    
    # Si plusieurs ateliers, ajouter le détail par atelier
    if len(ateliers_list) > 1:
        writer.writerow([])
        writer.writerow(['DÉTAIL PAR ATELIER'])
        writer.writerow([])
        writer.writerow([
            'Atelier', 'PPM Officiel', 'PPM Non Officiel', 'Taux Retouche (%)', 
            'Taux Rebut (%)', 'Coût Rebut (TND)', 'CNQ (TND)', 'Réclamations'
        ])
        
        for atelier_detail in ateliers_list:
            detail_params = [mois, annee, atelier_detail]
            
            # Calculer les données pour chaque atelier
            pieces_exp_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_exportees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_recl_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_reclamees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            alertes_detail = conn.execute('''
                SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                FROM quantite_alertes 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_ret_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_retouchees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_reb_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_rebutees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            pieces_ctrl_detail = conn.execute('''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_controlees 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            cout_reb_detail = conn.execute('''
                SELECT COALESCE(SUM(montant), 0) as total 
                FROM cout_rebut 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            cnq_detail = conn.execute('''
                SELECT COALESCE(AVG(valeur), 0) as avg 
                FROM cnq 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['avg']
            
            recl_off_detail = conn.execute('''
                SELECT COALESCE(SUM(nombre), 0) as total 
                FROM reclamations_officielles 
                WHERE mois = ? AND annee = ? AND atelier = ?
            ''', detail_params).fetchone()['total']
            
            # Calculer les indicateurs pour cet atelier
            ppm_off_detail = (pieces_recl_detail / pieces_exp_detail * 1000000) if pieces_exp_detail > 0 else 0
            ppm_non_off_detail = (alertes_detail / pieces_exp_detail * 1000000) if pieces_exp_detail > 0 else 0
            taux_ret_detail = (pieces_ret_detail / pieces_ctrl_detail * 100) if pieces_ctrl_detail > 0 else 0
            taux_reb_detail = (pieces_reb_detail / pieces_ctrl_detail * 100) if pieces_ctrl_detail > 0 else 0
            
            writer.writerow([
                atelier_detail, 
                round(ppm_off_detail, 2), 
                round(ppm_non_off_detail, 2), 
                round(taux_ret_detail, 2), 
                round(taux_reb_detail, 2), 
                round(cout_reb_detail, 2), 
                round(cnq_detail, 2), 
                recl_off_detail
            ])
    
    conn.close()
    
    # Préparer le fichier pour téléchargement
    output.seek(0)
    mem = io.BytesIO()
    mem.write(output.getvalue().encode('utf-8'))
    mem.seek(0)
    
    filename = f'consultation_qualite_{atelier.replace(" ", "_")}_{month_names[mois]}_{annee}_{datetime.now().strftime("%Y%m%d_%H%M%S")}.csv'
    
    return send_file(
        mem,
        as_attachment=True,
        download_name=filename,
        mimetype='text/csv'
    )
# Ajoutez cette route à votre app.py, par exemple après la route dashboard

@app.route('/quality/data_entry')
@login_required
def quality_data_entry():
    """Page de saisie des données qualité"""
    if session.get('user_role') not in ['Administration', 'Chargé Qualité/Animateur Qualité']:
        flash('Accès non autorisé', 'error')
        return redirect(url_for('dashboard'))
    
    return render_template('quality_manager/data_entry.html')
# Ajoutez ces nouvelles routes à votre app.py existant



@app.route('/api/bulk_chart_data')
@login_required
def get_bulk_chart_data():
    """Récupérer les données pour tous les indicateurs d'un atelier"""
    atelier = request.args.get('atelier')
    uap = request.args.get('uap', 'all')
    
    if not atelier:
        return jsonify({'error': 'Atelier requis'}), 400
    
    current_year = datetime.now().year
    previous_year = current_year - 1
    
    conn = get_db_connection()
    
    # Construire la condition UAP
    uap_condition = ""
    uap_params = []
    if uap != 'all':
        uap_condition = " AND uap = ?"
        uap_params = [uap]
    
    # Récupérer toutes les données mensuelles pour les deux années
    result = {
        'atelier': atelier,
        'uap': uap,
        'current_year': current_year,
        'previous_year': previous_year,
        'indicators': {}
    }
    
    # Liste des indicateurs à traiter
    indicators = [
        'ppm_officiel', 'ppm_non_officiel', 'cnq', 'cout_rebut',
        'taux_retouche', 'taux_rebut', 'nombre_reclamations'
    ]
    
    for indicator in indicators:
        current_data = []
        previous_data = []
        
        for mois in range(1, 13):
            # Calculer pour l'année actuelle
            current_value = calculate_monthly_indicator(conn, indicator, mois, current_year, atelier, uap_condition, uap_params)
            previous_value = calculate_monthly_indicator(conn, indicator, mois, previous_year, atelier, uap_condition, uap_params)
            
            current_data.append(round(current_value, 2))
            previous_data.append(round(previous_value, 2))
        
        # Récupérer le seuil
        indicateur_mapping = {
            'ppm_officiel': 'PPM Officiel',
            'ppm_non_officiel': 'PPM Non Officiel',
            'cnq': 'CNQ',
            'cout_rebut': 'Coût de Rebut',
            'taux_retouche': 'Taux de Retouche',
            'taux_rebut': 'Taux de Rebut',
            'nombre_reclamations': 'Nombre de Réclamations'
        }
        
        indicateur_name = indicateur_mapping.get(indicator, indicator)
        seuil_data = conn.execute('''
            SELECT valeur_seuil FROM seuils 
            WHERE atelier = ? AND indicateur = ?
            ORDER BY date_modification DESC 
            LIMIT 1
        ''', (atelier, indicateur_name)).fetchone()
        
        seuil = seuil_data['valeur_seuil'] if seuil_data else get_default_threshold(indicator)
        
        result['indicators'][indicator] = {
            'name': indicateur_name,
            'current_data': current_data,
            'previous_data': previous_data,
            'threshold': seuil
        }
    
    conn.close()
    return jsonify(result)

def calculate_monthly_indicator(conn, indicator, mois, annee, atelier, uap_condition, uap_params):
    """Fonction utilitaire pour calculer un indicateur pour un mois donné"""
    base_params = [mois, annee, atelier] + uap_params
    
    if indicator in ['ppm_officiel', 'ppm_non_officiel']:
        pieces_exportees = conn.execute(f'''
            SELECT COALESCE(SUM(quantite), 0) as total 
            FROM pieces_exportees 
            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
        ''', base_params).fetchone()['total'] or 1
        
        if indicator == 'ppm_officiel':
            pieces_reclamees = conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_reclamees 
                WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
            ''', base_params).fetchone()['total']
            return (pieces_reclamees / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
        else:  # ppm_non_officiel
            alertes = conn.execute(f'''
                SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                FROM quantite_alertes 
                WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
            ''', base_params).fetchone()['total']
            return (alertes / pieces_exportees * 1000000) if pieces_exportees > 0 else 0
            
    elif indicator in ['taux_retouche', 'taux_rebut']:
        pieces_controlees = conn.execute(f'''
            SELECT COALESCE(SUM(quantite), 0) as total 
            FROM pieces_controlees 
            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
        ''', base_params).fetchone()['total'] or 1
        
        if indicator == 'taux_retouche':
            pieces_retouchees = conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_retouchees 
                WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
            ''', base_params).fetchone()['total']
            return (pieces_retouchees / pieces_controlees * 100) if pieces_controlees > 0 else 0
        else:  # taux_rebut
            pieces_rebutees = conn.execute(f'''
                SELECT COALESCE(SUM(quantite), 0) as total 
                FROM pieces_rebutees 
                WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
            ''', base_params).fetchone()['total']
            return (pieces_rebutees / pieces_controlees * 100) if pieces_controlees > 0 else 0
            
    elif indicator == 'cout_rebut':
        cout = conn.execute(f'''
            SELECT COALESCE(SUM(montant), 0) as total 
            FROM cout_rebut 
            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
        ''', base_params).fetchone()['total']
        return cout
        
    elif indicator == 'cnq':
        cnq = conn.execute(f'''
            SELECT COALESCE(AVG(valeur), 0) as avg 
            FROM cnq 
            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
        ''', base_params).fetchone()['avg']
        return cnq
        
    elif indicator == 'nombre_reclamations':
        reclamations = conn.execute(f'''
            SELECT COALESCE(SUM(nombre), 0) as total 
            FROM reclamations_officielles 
            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
        ''', base_params).fetchone()['total']
        return reclamations
    
    return 0

def get_default_threshold(indicator):
    """Retourner les seuils par défaut pour les indicateurs"""
    default_thresholds = {
        'ppm_officiel': 100,
        'ppm_non_officiel': 500,
        'taux_rebut': 2.0,
        'taux_retouche': 5.0,
        'cnq': 1000,
        'cout_rebut': 500,
        'nombre_reclamations': 10
    }
    return default_thresholds.get(indicator, 0)

@app.route('/api/chart_export/<indicator>')
@login_required
def export_chart_data(indicator):
    """Exporter les données d'un graphique spécifique en CSV"""
    atelier = request.args.get('atelier')
    uap = request.args.get('uap', 'all')
    
    if not atelier:
        return jsonify({'error': 'Atelier requis'}), 400
    
    # Récupérer les données du graphique
    response = get_chart_analysis(indicator)
    if response.status_code != 200:
        return response
    
    data = response.get_json()
    
    # Créer le CSV
    output = io.StringIO()
    writer = csv.writer(output)
    
    # En-têtes
    writer.writerow([
        f'Données {data["indicator_name"]} - {atelier}',
        f'UAP: {uap}',
        f'Seuil: {data["threshold"]}'
    ])
    writer.writerow([])
    writer.writerow(['Mois', f'{data["current_year"]}', f'{data["previous_year"]}', 'Seuil'])
    
    # Données mensuelles
    months = ['Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
              'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre']
    
    for i in range(12):
        writer.writerow([
            months[i],
            data['current_data'][i],
            data['previous_data'][i],
            data['threshold']
        ])
    
    # Statistiques d'analyse
    writer.writerow([])
    writer.writerow(['Analyse'])
    writer.writerow(['Moyenne année actuelle', data['analysis']['current_average']])
    writer.writerow(['Moyenne année précédente', data['analysis']['previous_average']])
    writer.writerow(['Amélioration (%)', data['analysis']['improvement_percentage']])
    writer.writerow(['Tendance', data['analysis']['trend']])
    writer.writerow(['Dépassements seuil (actuel)', data['analysis']['current_threshold_violations']])
    writer.writerow(['Dépassements seuil (précédent)', data['analysis']['previous_threshold_violations']])
    
    # Préparer le fichier pour téléchargement
    output.seek(0)
    mem = io.BytesIO()
    mem.write(output.getvalue().encode('utf-8'))
    mem.seek(0)
    
    filename = f'{indicator}_{atelier.replace(" ", "_")}_{datetime.now().strftime("%Y%m%d_%H%M%S")}.csv'
    
    return send_file(
        mem,
        as_attachment=True,
        download_name=filename,
        mimetype='text/csv'
    )

@app.route('/api/get_threshold_by_indicator')
@login_required
def get_threshold_by_indicator():
    """Obtenir le seuil pour un indicateur et atelier spécifique"""
    atelier = request.args.get('atelier')
    indicateur = request.args.get('indicateur')
    
    if not all([atelier, indicateur]):
        return jsonify({'error': 'Atelier et indicateur requis'}), 400
    
    conn = get_db_connection()
    
    # Récupérer le seuil le plus récent
    seuil = conn.execute('''
        SELECT valeur_seuil FROM seuils 
        WHERE atelier = ? AND indicateur = ?
        ORDER BY date_modification DESC 
        LIMIT 1
    ''', (atelier, indicateur)).fetchone()
    
    conn.close()
    
    if seuil:
        return jsonify({'seuil': seuil['valeur_seuil']})
    else:
        # Retourner des seuils par défaut si non trouvé
        default_thresholds = {
            'PPM Officiel': 100,
            'PPM Non Officiel': 500,
            'Taux de Rebut': 2.0,
            'Taux de Retouche': 5.0,
            'CNQ': 1000,
            'Coût de Rebut': 500,
            'Nombre de Réclamations': 10
        }
        return jsonify({'seuil': default_thresholds.get(indicateur, 0)})

@app.route('/api/chart_analysis/<indicator>')
@login_required  
def get_chart_analysis(indicator):
    """Analyser un indicateur spécifique avec comparaison temporelle"""
    atelier = request.args.get('atelier')
    uap = request.args.get('uap', 'all')
    
    if not atelier:
        return jsonify({'error': 'Atelier requis'}), 400
    
    # Valider l'indicateur
    valid_indicators = ['ppm_officiel', 'ppm_non_officiel', 'cnq', 'cout_rebut', 
                       'taux_retouche', 'taux_rebut', 'nombre_reclamations']
    
    if indicator not in valid_indicators:
        return jsonify({'error': 'Indicateur non valide'}), 400
    
    current_year = datetime.now().year
    previous_year = current_year - 1
    
    conn = get_db_connection()
    
    # Construire la condition UAP
    uap_condition = ""
    uap_params = []
    if uap != 'all':
        uap_condition = " AND uap = ?"
        uap_params = [uap]
    
    # Calculer les données pour les 12 mois des deux années
    current_data = []
    previous_data = []
    
    for mois in range(1, 13):
        current_value = calculate_monthly_indicator(conn, indicator, mois, current_year, atelier, uap_condition, uap_params)
        previous_value = calculate_monthly_indicator(conn, indicator, mois, previous_year, atelier, uap_condition, uap_params)
        
        current_data.append(round(current_value, 2))
        previous_data.append(round(previous_value, 2))
    
    # Récupérer le seuil
    indicateur_mapping = {
        'ppm_officiel': 'PPM Officiel',
        'ppm_non_officiel': 'PPM Non Officiel',
        'cnq': 'CNQ',
        'cout_rebut': 'Coût de Rebut',
        'taux_retouche': 'Taux de Retouche',
        'taux_rebut': 'Taux de Rebut',
        'nombre_reclamations': 'Nombre de Réclamations'
    }
    
    indicateur_name = indicateur_mapping.get(indicator, indicator)
    
    seuil_data = conn.execute('''
        SELECT valeur_seuil FROM seuils 
        WHERE atelier = ? AND indicateur = ?
        ORDER BY date_modification DESC 
        LIMIT 1
    ''', (atelier, indicateur_name)).fetchone()
    
    seuil = seuil_data['valeur_seuil'] if seuil_data else get_default_threshold(indicator)
    
    # Calculer des statistiques d'analyse
    current_avg = sum(current_data) / len(current_data) if current_data else 0
    previous_avg = sum(previous_data) / len(previous_data) if previous_data else 0
    improvement = ((previous_avg - current_avg) / previous_avg * 100) if previous_avg > 0 else 0
    
    # Compter les dépassements de seuil
    current_threshold_violations = sum(1 for value in current_data if value > seuil) if seuil > 0 else 0
    previous_threshold_violations = sum(1 for value in previous_data if value > seuil) if seuil > 0 else 0
    
    conn.close()
    
    return jsonify({
        'indicator': indicator,
        'indicator_name': indicateur_name,
        'atelier': atelier,
        'uap': uap,
        'current_year': current_year,
        'previous_year': previous_year,
        'current_data': current_data,
        'previous_data': previous_data,
        'threshold': seuil,
        'analysis': {
            'current_average': round(current_avg, 2),
            'previous_average': round(previous_avg, 2),
            'improvement_percentage': round(improvement, 2),
            'current_threshold_violations': current_threshold_violations,
            'previous_threshold_violations': previous_threshold_violations,
            'trend': 'amélioration' if improvement > 0 else 'dégradation' if improvement < 0 else 'stable'
        }
    })





@app.route('/api/search_data')
def search_data():
    """API endpoint pour rechercher les données selon les critères"""
    try:
        print("=== DEBUG: Début de search_data ===")
        
        # Récupération des paramètres
        atelier = request.args.get('atelier')
        mois = request.args.get('mois')
        annee = request.args.get('annee')
        
        print(f"Paramètres reçus: atelier={atelier}, mois={mois}, annee={annee}")

        if not all([atelier, mois, annee]):
            print("Erreur: Paramètres manquants")
            return jsonify({'error': 'Paramètres manquants'}), 400

        # Convertir en int avec gestion d'erreur
        try:
            mois = int(mois)
            annee = int(annee)
            print(f"Paramètres convertis: mois={mois}, annee={annee}")
        except ValueError as e:
            print(f"Erreur de conversion: {e}")
            return jsonify({'error': 'Paramètres invalides'}), 400

        # Connexion à la base de données
        try:
            conn = sqlite3.connect('qualite_ADI_final.sqlite')
            print("Connexion à la base réussie")
        except Exception as e:
            print(f"Erreur de connexion à la base: {e}")
            return jsonify({'error': f'Erreur de base de données: {str(e)}'}), 500

        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()

        # DIAGNOSTIC: Vérifier la structure de la base
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
        tables = [row[0] for row in cursor.fetchall()]
        print(f"Tables disponibles: {tables}")

        # DIAGNOSTIC: Vérifier un échantillon de données
        if 'pieces_exportees' in tables:
            cursor.execute("SELECT * FROM pieces_exportees LIMIT 5")
            sample = cursor.fetchall()
            print(f"Échantillon pieces_exportees: {[dict(row) for row in sample]}")
            
            # Vérifier les colonnes
            cursor.execute("PRAGMA table_info(pieces_exportees)")
            columns = cursor.fetchall()
            print(f"Colonnes pieces_exportees: {[col[1] for col in columns]}")

        # Définir les indicateurs avec leurs tables et champs
        indicators = [
            ('pieces_retouchees', 'quantite'),
            ('pieces_rebutees', 'quantite'),
            ('pieces_controlees', 'quantite'),
            ('cout_rebut', 'montant'),
            ('cnq', 'valeur'),
            ('pieces_reclamees', 'quantite'),
            ('pieces_exportees', 'quantite'),
            ('quantite_alertes', 'nombre_alertes'),
            ('reclamations_officielles', 'nombre')
        ]

        results = {}
        
        # Déterminer le type de requête selon l'atelier sélectionné
        is_aggregate = atelier in ['UAP1', 'UAP2', 'Total']

        for table, field in indicators:
            try:
                print(f"=== Traitement de la table: {table} ===")
                
                # Vérifier que la table existe
                if table not in tables:
                    print(f"Table {table} n'existe pas, passage à la suivante")
                    continue

                # Vérifier les colonnes de la table
                cursor.execute(f"PRAGMA table_info({table})")
                columns_info = cursor.fetchall()
                column_names = [col[1] for col in columns_info]
                print(f"Colonnes de {table}: {column_names}")
                
                # Vérifier si les colonnes nécessaires existent
                required_columns = ['mois', 'annee', field]
                if not is_aggregate or table != 'cnq':
                    required_columns.append('atelier')
                required_columns.append('uap')
                missing_cols = [col for col in required_columns if col not in column_names]
                if missing_cols:
                    print(f"Colonnes manquantes dans {table}: {missing_cols}")

                # Spécial pour cnq: seulement disponible pour atelier='Total'
                if table == 'cnq' and atelier != 'Total':
                    print(f"CNQ non disponible pour {atelier}, définition à 0")
                    results[table] = {'total': 0.0, 'type': 'aggregate'} if is_aggregate else []
                    continue

                # Compter le total d'enregistrements pour debug
                cursor.execute(f"SELECT COUNT(*) as count FROM {table}")
                total_count = cursor.fetchone()
                print(f"Total d'enregistrements dans {table}: {total_count['count']}")

                # Compter pour la période spécifique
                if 'mois' in column_names and 'annee' in column_names:
                    cursor.execute(f"SELECT COUNT(*) as count FROM {table} WHERE mois = ? AND annee = ?", 
                                 (mois, annee))
                    period_count = cursor.fetchone()
                    print(f"Enregistrements pour {mois}/{annee} dans {table}: {period_count['count']}")
                else:
                    print(f"Colonnes mois/annee manquantes dans {table}")
                    continue

                # Exécuter la requête appropriée
                if is_aggregate:
                    if table == 'cnq' and atelier == 'Total':
                        query = f"""
                        SELECT SUM(CASE WHEN {field} IS NOT NULL THEN CAST({field} AS REAL) ELSE 0 END) as total,
                               MAX(CASE WHEN seuil IS NOT NULL THEN CAST(seuil AS REAL) ELSE 0 END) as seuil
                        FROM {table} 
                        WHERE mois = ? AND annee = ?
                        """
                        cursor.execute(query, (mois, annee))
                        result = cursor.fetchone()
                        total_value = float(result['total']) if result['total'] is not None else 0.0
                        seuil_value = float(result['seuil']) if result['seuil'] is not None else 0.0
                        results[table] = {
                            'total': total_value,
                            'seuil': seuil_value,
                            'type': 'aggregate'
                        }
                        print(f"Résultat agrégé pour {table}: total={total_value}, seuil={seuil_value}")
                    else:
                        if atelier in ['UAP1', 'UAP2'] and 'uap' in column_names:
                            # Utiliser la colonne uap pour UAP1/UAP2
                            query = f"""
                            SELECT SUM(CASE WHEN {field} IS NOT NULL THEN CAST({field} AS REAL) ELSE 0 END) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ? AND uap = ?
                            """
                            cursor.execute(query, (mois, annee, atelier))
                        elif atelier == 'Total':
                            # Pour Total (Usine Complète), sommer tout
                            query = f"""
                            SELECT SUM(CASE WHEN {field} IS NOT NULL THEN CAST({field} AS REAL) ELSE 0 END) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ?
                            """
                            cursor.execute(query, (mois, annee))
                        else:
                            # Alternative si pas de uap: sommer par ateliers spécifiques
                            if atelier == 'UAP1':
                                ateliers_list = "('Protections thermiques', 'Isolant souple')"
                            else:  # UAP2
                                ateliers_list = "('Manchons', 'Colliers', 'Racks', 'Moulage', 'Composite', 'Système de visualisation')"
                            
                            query = f"""
                            SELECT SUM(CASE WHEN {field} IS NOT NULL THEN CAST({field} AS REAL) ELSE 0 END) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ? AND atelier IN {ateliers_list}
                            """
                            cursor.execute(query, (mois, annee))
                        
                        result = cursor.fetchone()
                        total_value = result['total'] if result and result['total'] is not None else 0
                        print(f"Résultat agrégé pour {table}: {total_value}")
                        
                        results[table] = {
                            'total': float(total_value),
                            'type': 'aggregate'
                        }
                else:
                    # Pour un atelier spécifique
                    if 'atelier' in column_names:
                        query = f"""
                        SELECT * FROM {table} 
                        WHERE mois = ? AND annee = ? AND atelier = ?
                        """
                        cursor.execute(query, (mois, annee, atelier))
                    else:
                        # Si pas de colonne atelier, récupérer toutes les données de la période
                        query = f"""
                        SELECT * FROM {table} 
                        WHERE mois = ? AND annee = ?
                        """
                        cursor.execute(query, (mois, annee))
                    
                    rows = cursor.fetchall()
                    print(f"Nombre de lignes trouvées pour {table}: {len(rows)}")
                    
                    if rows:
                        results[table] = [dict(row) for row in rows]
                        print(f"Premier enregistrement: {dict(rows[0])}")
                    else:
                        results[table] = []

            except Exception as e:
                print(f"Erreur lors du traitement de la table {table}: {e}")
                traceback.print_exc()
                results[table] = [] if not is_aggregate else {'total': 0.0, 'type': 'aggregate'}
                continue

        # Récupérer les seuils depuis la table seuils pour l'atelier spécifié
        try:
            if atelier in ['UAP1', 'UAP2', 'Total']:
                # Pour les agrégations (UAP1, UAP2, Total), utiliser le seuil correspondant
                query = """
                SELECT indicateur, valeur_seuil 
                FROM seuils 
                WHERE atelier = ?
                """
                cursor.execute(query, (atelier,))
            else:
                # Pour les ateliers spécifiques
                query = """
                SELECT indicateur, valeur_seuil 
                FROM seuils 
                WHERE atelier = ?
                """
                cursor.execute(query, (atelier,))
            
            seuils_rows = cursor.fetchall()
            seuils = {row['indicateur']: float(row['valeur_seuil']) if row['valeur_seuil'] is not None else 0.0 for row in seuils_rows}
            print(f"Seuils récupérés pour atelier {atelier}: {seuils}")
            
            # Ajouter le seuil pour CNQ depuis la table cnq (inchangé)
            if 'cnq' in results and 'seuil' in results['cnq']:
                seuils['CNQ'] = results['cnq']['seuil']
                print(f"Seuil CNQ ajouté: {seuils['CNQ']}")
            
            results['seuils'] = seuils
        except Exception as e:
            print(f"Erreur lors de la récupération des seuils: {e}")
            seuils = {
                'PPM Officiel': 100.0,
                'PPM Non Officiel': 500.0,
                'Taux de Retouche': 0.0,
                'Taux de Rebut': 2.0,
                'Coût de Rebut': 500.0,
                'Nombre de reclamation': 5.0,
                'CNQ': 1000.0 if atelier == 'Total' else 0.0
            }
            print(f"Utilisation des seuils par défaut: {seuils}")
            results['seuils'] = seuils

        conn.close()
        
        # Ajouter les informations sur le type de requête
        results['query_info'] = {
            'atelier': atelier,
            'mois': mois,
            'annee': annee,
            'is_aggregate': is_aggregate
        }
        
        # DIAGNOSTIC FINAL
        print("=== DIAGNOSTIC FINAL ===")
        for key, value in results.items():
            if key != 'query_info' and key != 'seuils':
                if isinstance(value, dict) and 'total' in value:
                    print(f"{key}: {value['total']}")
                elif isinstance(value, list):
                    print(f"{key}: {len(value)} enregistrements")
        print(f"seuils: {results['seuils']}")
        
        return jsonify(results)

    except Exception as e:
        print(f"Erreur générale dans search_data: {e}")
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/export_data')
def export_data():
    """API endpoint pour exporter les données en Excel"""
    try:
        atelier = request.args.get('atelier')
        mois = request.args.get('mois')
        annee = request.args.get('annee')

        if not all([atelier, mois, annee]):
            return jsonify({'error': 'Paramètres manquants'}), 400

        mois = int(mois)
        annee = int(annee)

        conn = sqlite3.connect('qualite_ADI_final.sqlite')
        
        # Noms des mois en français
        month_names = ['', 'Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
                      'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre']

        # Définir les indicateurs avec leurs noms
        indicators = [
            ('pieces_retouchees', 'Pièces Retouchées', 'quantite'),
            ('pieces_rebutees', 'Pièces Rebutées', 'quantite'),
            ('pieces_controlees', 'Pièces Contrôlées', 'quantite'),
            ('cout_rebut', 'Coût du Rebut', 'montant'),
            ('cnq', 'CNQ', 'valeur'),
            ('pieces_reclamees', 'Pièces Réclamées', 'quantite'),
            ('pieces_exportees', 'Pièces Exportées', 'quantite'),
            ('quantite_alertes', 'Quantité d\'Alertes', 'nombre_alertes'),
            ('reclamations_officielles', 'Réclamations Officielles', 'nombre')
        ]

        # Créer un buffer Excel en mémoire
        output = io.BytesIO()
        
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            
            is_aggregate = atelier in ['UAP1', 'UAP2', 'Usine Complète']
            
            if is_aggregate:
                # Feuille de synthèse agrégée
                synthesis_data = []
                
                for table, indicator_name, field in indicators:
                    try:
                        if atelier == 'UAP1':
                            query = f"""
                            SELECT SUM(CAST({field} AS REAL)) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ? 
                            AND (atelier IN ('Isolation Thermique', 'Isolation Souple') OR uap = 'UAP1')
                            """
                        elif atelier == 'UAP2':
                            query = f"""
                            SELECT SUM(CAST({field} AS REAL)) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ? 
                            AND (atelier NOT IN ('Isolation Thermique', 'Isolation Souple') OR uap = 'UAP2')
                            """
                        else:  # Usine Complète
                            query = f"""
                            SELECT SUM(CAST({field} AS REAL)) as total
                            FROM {table} 
                            WHERE mois = ? AND annee = ?
                            """
                        
                        df = pd.read_sql_query(query, conn, params=(mois, annee))
                        
                        if not df.empty and df['total'].iloc[0] is not None:
                            total_value = float(df['total'].iloc[0])
                            synthesis_data.append({
                                'Indicateur': indicator_name,
                                'Valeur': total_value
                            })
                    except Exception as e:
                        print(f"Erreur pour {table}: {e}")
                        continue

                if synthesis_data:
                    synthesis_df = pd.DataFrame(synthesis_data)
                    sheet_name = f'Synthèse {atelier}'
                    synthesis_df.to_excel(writer, sheet_name=sheet_name, index=False)
                    
            else:
                # Feuille détaillée pour un atelier spécifique
                synthesis_data = []
                
                for table, indicator_name, field in indicators:
                    try:
                        query = f"""
                        SELECT * FROM {table} 
                        WHERE mois = ? AND annee = ? AND atelier = ?
                        """
                        df_detail = pd.read_sql_query(query, conn, params=(mois, annee, atelier))

                        if not df_detail.empty:
                            # Ajouter à la synthèse
                            total_value = df_detail[field].sum()
                            synthesis_data.append({
                                'Indicateur': indicator_name,
                                'Valeur': total_value
                            })
                            
                            # Créer une feuille détaillée
                            sheet_name = indicator_name[:31]  # Limite Excel
                            df_detail.to_excel(writer, sheet_name=sheet_name, index=False)
                    except Exception as e:
                        print(f"Erreur pour {table}: {e}")
                        continue

                # Créer la feuille de synthèse
                if synthesis_data:
                    synthesis_df = pd.DataFrame(synthesis_data)
                    synthesis_df.to_excel(writer, sheet_name='Synthèse', index=False)

        conn.close()
        
        # Préparer le fichier pour le téléchargement
        output.seek(0)
        
        # Nom du fichier
        filename = f"donnees_qualite_{atelier.replace(' ', '_')}_{month_names[mois]}_{annee}.xlsx"
        
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        print(f"Erreur dans export_data: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/consultation')
def consultation():
    """Route pour afficher la page de consultation"""
    return render_template('consultation.html')

# Fonction utilitaire pour calculer les totaux par UAP
def calculate_uap_totals(data, field_name):
    """
    Calcule les totaux par UAP à partir des données
    
    Args:
        data: Liste de dictionnaires contenant les données
        field_name: Nom du champ à sommer
    
    Returns:
        dict: Dictionnaire avec les totaux UAP1, UAP2 et total
    """
    uap1_total = 0
    uap2_total = 0
    total = 0
    
    for row in data:
        value = float(row.get(field_name, 0) or 0)
        uap = str(row.get('uap', '')).upper()
        
        if '1' in uap or uap == 'UAP1':
            uap1_total += value
        elif '2' in uap or uap == 'UAP2':
            uap2_total += value
        
        total += value
    
    return {
        'uap1': uap1_total,
        'uap2': uap2_total,
        'total': total
    }

# Fonction pour obtenir les seuils d'un atelier
def get_seuils_for_atelier(atelier):
    """
    Récupère les seuils définis pour un atelier donné
    
    Args:
        atelier: Nom de l'atelier
    
    Returns:
        dict: Dictionnaire des seuils par indicateur
    """
    try:
        conn = sqlite3.connect('adi_quality.db')
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT indicateur, valeur_seuil 
            FROM seuils 
            WHERE atelier = ? OR atelier = 'Tous'
            ORDER BY date_modification DESC
        """, (atelier,))
        
        seuils = {}
        for row in cursor.fetchall():
            seuils[row['indicateur']] = row['valeur_seuil']
        
        conn.close()
        return seuils
        
    except Exception as e:
        print(f"Erreur lors de la récupération des seuils: {e}")
        return {}

# API pour obtenir les seuils
@app.route('/api/get_seuils/<atelier>')
def get_seuils(atelier):
    """API endpoint pour récupérer les seuils d'un atelier"""
    try:
        seuils = get_seuils_for_atelier(atelier)
        return jsonify(seuils)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Fonction pour comparer les valeurs aux seuils
def compare_with_thresholds(value, threshold, indicator_type='quality'):
    """
    Compare une valeur avec son seuil et retourne le statut
    
    Args:
        value: Valeur à comparer
        threshold: Seuil de référence
        indicator_type: Type d'indicateur ('quality' ou 'quantity')
    
    Returns:
        str: 'success', 'warning', ou 'danger'
    """
    if threshold is None or threshold == 0:
        return 'info'
    
    ratio = value / threshold
    
    if indicator_type == 'quality':
        # Pour les indicateurs de qualité, moins c'est mieux
        if ratio <= 0.8:
            return 'success'
        elif ratio <= 1.0:
            return 'warning'
        else:
            return 'danger'
    else:
        # Pour les indicateurs de quantité, plus c'est mieux
        if ratio >= 1.2:
            return 'success'
        elif ratio >= 0.8:
            return 'warning'
        else:
            return 'danger'

# API pour obtenir les données avec comparaison aux seuils
@app.route('/api/search_data_with_thresholds')
def search_data_with_thresholds():
    """API endpoint pour rechercher les données avec comparaison aux seuils"""
    try:
        atelier = request.args.get('atelier')
        mois = request.args.get('mois')
        annee = request.args.get('annee')

        if not all([atelier, mois, annee]):
            return jsonify({'error': 'Paramètres manquants'}), 400

        # Récupérer les données
        search_response = search_data()
        if search_response.status_code != 200:
            return search_response
        
        data = search_response.get_json()
        
        # Récupérer les seuils
        seuils = get_seuils_for_atelier(atelier)
        
        # Enrichir les données avec les comparaisons de seuils
        enriched_data = {}
        
        indicator_mapping = {
            'pieces_retouchees': ('Pièces Retouchées', 'quality', 'quantite'),
            'pieces_rebutees': ('Pièces Rebutées', 'quality', 'quantite'),
            'pieces_controlees': ('Pièces Contrôlées', 'quantity', 'quantite'),
            'cout_rebut': ('Coût du Rebut', 'quality', 'montant'),
            'cnq': ('CNQ', 'quality', 'valeur'),
            'pieces_reclamees': ('Pièces Réclamées', 'quality', 'quantite'),
            'pieces_exportees': ('Pièces Exportées', 'quantity', 'quantite'),
            'quantite_alertes': ('Quantité d\'Alertes', 'quality', 'nombre_alertes'),
            'reclamations_officielles': ('Réclamations Officielles', 'quality', 'nombre')
        }
        
        for table_name, table_data in data.items():
            if table_name in indicator_mapping:
                indicator_name, indicator_type, field_name = indicator_mapping[table_name]
                
                # Calculer les totaux
                totals = calculate_uap_totals(table_data, field_name)
                
                # Comparer avec le seuil
                threshold = seuils.get(indicator_name)
                status = 'info'
                if threshold:
                    status = compare_with_thresholds(totals['total'], threshold, indicator_type)
                
                enriched_data[table_name] = {
                    'data': table_data,
                    'totals': totals,
                    'threshold': threshold,
                    'status': status,
                    'indicator_name': indicator_name
                }
        
        return jsonify(enriched_data)

    except Exception as e:
        print(f"Erreur dans search_data_with_thresholds: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/admin/update_cnq', methods=['POST'])
@login_required
@admin_required
def update_cnq():
    """
    Met à jour la valeur du CNQ (Coût de Non-Qualité) pour l'usine
    """
    try:
        # Vérification de l'authentification et des permissions
        if 'user_id' not in session:
            flash('Vous devez être connecté pour accéder à cette page.', 'error')
            return redirect(url_for('login'))
        
        # Vérifier le rôle de l'utilisateur - CORRECTION ICI
        user_role = session.get('user_role')  # Changé de 'role' à 'user_role'
        print(f"[DEBUG] Rôle utilisateur: '{user_role}'")  # Pour déboguer
        
        # Liste des rôles autorisés (ajustez selon vos rôles réels)
        allowed_roles = [
            'Administration', 
            'Chargé Qualité/Animateur Qualité',
            'Admin',  # Ajout d'autres variantes possibles
            'Administrateur',
            'Chargé Qualité',
            'Animateur Qualité'
        ]
        
        if user_role not in allowed_roles:
            print(f"[DEBUG] Rôle '{user_role}' non autorisé. Rôles autorisés: {allowed_roles}")
            flash(f'Votre rôle "{user_role}" n\'a pas les permissions nécessaires pour modifier les données CNQ.', 'error')
            return redirect(url_for('dashboard'))
        
        # Récupération des données du formulaire
        cnq_value = request.form.get('cnq_value')
        cnq_month = request.form.get('cnq_month')
        cnq_year = request.form.get('cnq_year')
        cnq_threshold_value = request.form.get('cnq_threshold_value')  # Pour modification du seuil
        
        # Validation des données obligatoires pour ajout/modification CNQ
        if cnq_value and cnq_month and cnq_year:
            try:
                cnq_value = float(cnq_value)
                cnq_month = int(cnq_month)
                cnq_year = int(cnq_year)
                
                # Validation des plages de valeurs
                if cnq_value < 0 or cnq_value > 100:
                    flash('La valeur CNQ doit être comprise entre 0 et 100%.', 'error')
                    return redirect(url_for('admin_thresholds'))
                
                if cnq_month < 1 or cnq_month > 12:
                    flash('Le mois doit être compris entre 1 et 12.', 'error')
                    return redirect(url_for('admin_thresholds'))
                
                if cnq_year < 2020 or cnq_year > 2030:
                    flash('L\'année doit être comprise entre 2020 et 2030.', 'error')
                    return redirect(url_for('admin_thresholds'))
                
            except ValueError:
                flash('Valeurs invalides. Veuillez vérifier vos saisies.', 'error')
                return redirect(url_for('admin_thresholds'))
        
        # Validation du seuil CNQ si fourni
        if cnq_threshold_value:
            try:
                cnq_threshold_value = float(cnq_threshold_value)
                if cnq_threshold_value < 0 or cnq_threshold_value > 100:
                    flash('Le seuil CNQ doit être compris entre 0 et 100%.', 'error')
                    return redirect(url_for('admin_thresholds'))
            except ValueError:
                flash('Valeur de seuil CNQ invalide.', 'error')
                return redirect(url_for('admin_thresholds'))
        
        # Connexion à la base de données
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # **CAS 1: Modification du seuil CNQ uniquement**
        if cnq_threshold_value and not (cnq_value and cnq_month and cnq_year):
            # Mettre à jour le seuil pour toutes les lignes
            cursor.execute('''
                UPDATE cnq 
                SET seuil = ?, created_at = CURRENT_TIMESTAMP
                WHERE atelier = 'Total'
            ''', (cnq_threshold_value,))
            
            flash(f'Seuil CNQ mis à jour : {cnq_threshold_value}% pour toutes les données', 'success')
        
        # **CAS 2: Ajout/Modification valeur CNQ (avec ou sans nouveau seuil)**
        elif cnq_value and cnq_month and cnq_year:
            # Récupérer le seuil actuel ou utiliser le nouveau seuil
            if cnq_threshold_value:
                current_threshold = cnq_threshold_value
            else:
                # Récupérer le seuil existant
                cursor.execute('''
                    SELECT seuil FROM cnq 
                    WHERE atelier = 'Total' 
                    ORDER BY created_at DESC 
                    LIMIT 1
                ''')
                threshold_result = cursor.fetchone()
                current_threshold = threshold_result[0] if threshold_result else 15.0  # Valeur par défaut
            
            # Vérifier si une entrée existe déjà pour ce mois/année
            cursor.execute('''
                SELECT id FROM cnq 
                WHERE mois = ? AND annee = ? AND atelier = 'Total'
            ''', (cnq_month, cnq_year))
            
            existing_entry = cursor.fetchone()
            
            if existing_entry:
                # Mise à jour de la valeur CNQ existante
                cursor.execute('''
                    UPDATE cnq 
                    SET valeur = ?, seuil = ?, created_at = CURRENT_TIMESTAMP
                    WHERE mois = ? AND annee = ? AND atelier = 'Total'
                ''', (cnq_value, current_threshold, cnq_month, cnq_year))
                
                action_message = f'mise à jour pour {cnq_month:02d}/{cnq_year}'
            else:
                # Insertion d'une nouvelle entrée CNQ
                cursor.execute('''
                    INSERT INTO cnq (mois, annee, atelier, valeur, seuil)
                    VALUES (?, ?, 'Total', ?, ?)
                ''', (cnq_month, cnq_year, cnq_value, current_threshold))
                
                action_message = f'ajout pour {cnq_month:02d}/{cnq_year}'
            
            # Si un nouveau seuil est fourni, mettre à jour toutes les autres lignes aussi
            if cnq_threshold_value:
                cursor.execute('''
                    UPDATE cnq 
                    SET seuil = ? 
                    WHERE atelier = 'Total' AND NOT (mois = ? AND annee = ?)
                ''', (cnq_threshold_value, cnq_month, cnq_year))
            
            # Message de succès
            mois_names = {
                1: 'Janvier', 2: 'Février', 3: 'Mars', 4: 'Avril',
                5: 'Mai', 6: 'Juin', 7: 'Juillet', 8: 'Août',
                9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Décembre'
            }
            
            mois_name = mois_names.get(cnq_month, f'Mois {cnq_month}')
            
            success_msg = f'CNQ {action_message} : {cnq_value}% pour {mois_name} {cnq_year}'
            if cnq_threshold_value:
                success_msg += f' (Seuil mis à jour : {cnq_threshold_value}%)'
            
            flash(success_msg, 'success')
        
        # **CAS 3: Aucune donnée valide fournie**
        else:
            flash('Veuillez fournir soit une valeur CNQ (avec mois/année) soit un nouveau seuil CNQ.', 'error')
            return redirect(url_for('dashboard'))  # Changé temporairement
        
        # Commit des changements
        conn.commit()
        
        # Log de l'activité - CORRECTION ICI AUSSI
        user_name = session.get('user_name', 'Utilisateur inconnu')  # Changé de 'nom_compte' à 'user_name'
        if cnq_threshold_value and not (cnq_value and cnq_month and cnq_year):
            print(f"[CNQ THRESHOLD UPDATE] {user_name} a mis à jour le seuil CNQ: {cnq_threshold_value}%")
        elif cnq_value and cnq_month and cnq_year:
            print(f"[CNQ VALUE UPDATE] {user_name} a mis à jour CNQ: {cnq_value}% pour {cnq_month:02d}/{cnq_year}")
        
    except Exception as e:
        # Gestion des erreurs
        if 'conn' in locals():
            conn.rollback()
        
        print(f"Erreur lors de la mise à jour du CNQ : {str(e)}")
        flash('Une erreur s\'est produite lors de la mise à jour du CNQ. Veuillez réessayer.', 'error')
        
    finally:
        # Fermeture de la connexion
        if 'conn' in locals():
            conn.close()
    
    return redirect(url_for('admin_thresholds'))  # Changé temporairement - voir solutions ci-dessous



# API pour la comparaison temporelle - VERSION COMPLÈTE
@app.route('/api/charts_data_comparison')
@login_required
def get_charts_data_comparison():
    """API pour la comparaison entre deux années"""
    atelier = request.args.get('atelier')
    annee1 = request.args.get('annee1', type=int)
    annee2 = request.args.get('annee2', type=int)
    uap = request.args.get('uap', 'all')
    aggregate = request.args.get('aggregate')
    
    if not all([atelier, annee1, annee2]):
        return jsonify({'error': 'Paramètres manquants: atelier, annee1, annee2 requis'}), 400
    
    conn = get_db_connection()
    
    # Mêmes mappings que l'API principale
    atelier_mapping = {
        'Manchon': 'Manchons',
        'Collier de Fixation': 'Colliers',
        'Rack': 'Racks',
        'Moulage': 'Moulage',
        'Isolation Thermique': 'Protections thermiques',
        'Isolation Souple': 'Isolant souple',
        'Composite': 'Composite',
        'Atelier de Visualisation': 'Système de visualisation'
    }
    
    cout_rebut_mapping = {
        'Manchon': 'manchons',
        'Collier de Fixation': 'colliers',
        'Rack': 'rack',
        'Moulage': 'moulage',
        'Isolation Thermique': 'isolation termique',
        'Isolation Souple': 'isolation souple',
        'Composite': 'composite',
        'Atelier de Visualisation': 'isolation souple'
    }
    
    reclamations_mapping = {
        'Manchon': 'Manchons',
        'Collier de Fixation': 'Colliers',
        'Rack': 'Racks',
        'Moulage': 'Moulage',
        'Isolation Thermique': 'P.Thermiques',
        'Isolation Souple': 'I.Souples',
        'Composite': 'Composite',
        'Atelier de Visualisation': 'I.Souples'
    }
    
    seuil_mapping = {
        'ppm_officiel': 'PPM Officiel',
        'ppm_non_officiel': 'PPM Non Officiel', 
        'taux_rebut': 'Taux de Rebut',
        'taux_retouche': 'Taux de Retouche',
        'cout_rebut': 'Coût de Rebut',
        'cnq': 'CNQ',
        'nombre_reclamations': 'Nombre de reclamation'
    }
    
    try:
        # Déterminer les ateliers à traiter
        if atelier and atelier != 'Total':
            ateliers_interface = [a.strip() for a in atelier.split(',')]
        else:
            ateliers_interface = list(atelier_mapping.keys())

        # Construction des conditions UAP
        uap_condition = ""
        uap_params = []
        if uap != 'all':
            uap_condition = " AND uap = ?"
            uap_params = [uap]

        comparison_data = {}
        
        # Récupérer les données pour les deux années
        for annee in [annee1, annee2]:
            monthly_data = {}
            
            # Calculer les données mensuelles pour cette année
            for mois in range(1, 13):
                # Initialiser les totaux pour ce mois
                total_pieces_exportees = 0
                total_pieces_reclamees = 0
                total_alertes = 0
                total_pieces_retouchees = 0
                total_pieces_rebutees = 0
                total_pieces_controlees = 0
                total_cout_rebut = 0
                total_cnq_sum = 0
                cnq_count = 0
                total_reclamations_off = 0

                # Parcourir chaque atelier
                for atelier_interface in ateliers_interface:
                    # Obtenir le nom correct pour chaque table
                    atelier_standard = atelier_mapping.get(atelier_interface, atelier_interface)
                    atelier_cout = cout_rebut_mapping.get(atelier_interface, atelier_interface.lower())
                    atelier_recl = reclamations_mapping.get(atelier_interface, atelier_interface)
                    
                    base_params = [mois, annee, atelier_standard] + uap_params
                    
                    try:
                        # Pièces exportées
                        pieces_exp = conn.execute(f'''
                            SELECT COALESCE(SUM(quantite), 0) as total 
                            FROM pieces_exportees 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_pieces_exportees += pieces_exp
                        
                        # Pièces réclamées
                        pieces_recl = conn.execute(f'''
                            SELECT COALESCE(SUM(quantite), 0) as total 
                            FROM pieces_reclamees 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_pieces_reclamees += pieces_recl
                        
                        # Alertes
                        alertes = conn.execute(f'''
                            SELECT COALESCE(SUM(nombre_alertes), 0) as total 
                            FROM quantite_alertes 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_alertes += alertes
                        
                        # Pièces retouchées
                        pieces_ret = conn.execute(f'''
                            SELECT COALESCE(SUM(quantite), 0) as total 
                            FROM pieces_retouchees 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_pieces_retouchees += pieces_ret
                        
                        # Pièces rebutées
                        pieces_reb = conn.execute(f'''
                            SELECT COALESCE(SUM(quantite), 0) as total 
                            FROM pieces_rebutees 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_pieces_rebutees += pieces_reb
                        
                        # Pièces contrôlées
                        pieces_ctrl = conn.execute(f'''
                            SELECT COALESCE(SUM(quantite), 0) as total 
                            FROM pieces_controlees 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['total']
                        total_pieces_controlees += pieces_ctrl
                        
                        # Coût rebut
                        cout_params = [mois, annee, atelier_cout] + uap_params
                        cout = conn.execute(f'''
                            SELECT COALESCE(SUM(montant), 0) as total 
                            FROM cout_rebut 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', cout_params).fetchone()['total']
                        total_cout_rebut += cout
                        
                        # CNQ
                        cnq_val = conn.execute(f'''
                            SELECT COALESCE(AVG(valeur), 0) as avg 
                            FROM cnq 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', base_params).fetchone()['avg']
                        if cnq_val > 0:
                            total_cnq_sum += cnq_val
                            cnq_count += 1
                        
                        # Réclamations officielles
                        recl_params = [mois, annee, atelier_recl] + uap_params
                        recl_off = conn.execute(f'''
                            SELECT COALESCE(SUM(nombre), 0) as total 
                            FROM reclamations_officielles 
                            WHERE mois = ? AND annee = ? AND atelier = ?{uap_condition}
                        ''', recl_params).fetchone()['total']
                        total_reclamations_off += recl_off
                        
                    except Exception as e:
                        print(f"Erreur pour {atelier_interface}, mois {mois}, année {annee}: {e}")
                        continue

                # Calculer les indicateurs pour ce mois
                ppm_officiel = (total_pieces_reclamees / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
                ppm_non_officiel = (total_alertes / total_pieces_exportees * 1000000) if total_pieces_exportees > 0 else 0
                taux_retouche = (total_pieces_retouchees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0
                taux_rebut = (total_pieces_rebutees / total_pieces_controlees * 100) if total_pieces_controlees > 0 else 0
                cnq_moyenne = total_cnq_sum / cnq_count if cnq_count > 0 else 0

                monthly_data[mois] = {
                    'ppm_officiel': round(ppm_officiel, 2),
                    'ppm_non_officiel': round(ppm_non_officiel, 2),
                    'cout_rebut': round(total_cout_rebut, 2),
                    'nombre_reclamations': total_reclamations_off,
                    'cnq': round(cnq_moyenne, 2),
                    'taux_rebut': round(taux_rebut, 2),
                    'taux_retouche': round(taux_retouche, 2)
                }
            
            # Récupérer les seuils pour cette année (optionnel, généralement identiques)
            seuils = {}
            if len(ateliers_interface) > 1 or atelier == 'Total':
                # Pour plusieurs ateliers : moyenne des seuils
                for api_indicator, db_indicator in seuil_mapping.items():
                    try:
                        if len(ateliers_interface) == len(list(atelier_mapping.keys())):
                            # Total de tous les ateliers
                            seuil_result = conn.execute('''
                                SELECT AVG(valeur_seuil) as valeur_seuil
                                FROM seuils 
                                WHERE indicateur = ?
                            ''', [db_indicator]).fetchone()
                        else:
                            # Moyenne des ateliers sélectionnés
                            seuils_placeholders = ','.join(['?' for _ in ateliers_interface])
                            seuil_result = conn.execute(f'''
                                SELECT AVG(valeur_seuil) as valeur_seuil
                                FROM seuils 
                                WHERE atelier IN ({seuils_placeholders}) AND indicateur = ?
                            ''', ateliers_interface + [db_indicator]).fetchone()
                        
                        if seuil_result and seuil_result['valeur_seuil']:
                            seuils[api_indicator] = float(seuil_result['valeur_seuil'])
                        else:
                            seuils[api_indicator] = 0
                    except Exception as e:
                        print(f"Erreur seuil {api_indicator}: {e}")
                        seuils[api_indicator] = 0
            else:
                # Pour un seul atelier
                atelier_interface = ateliers_interface[0]
                
                for api_indicator, db_indicator in seuil_mapping.items():
                    try:
                        seuil_result = conn.execute('''
                            SELECT valeur_seuil 
                            FROM seuils
                            WHERE atelier = ? AND indicateur = ?
                            ORDER BY date_modification DESC
                            LIMIT 1
                        ''', (atelier_interface, db_indicator)).fetchone()
                        
                        if seuil_result:
                            seuils[api_indicator] = float(seuil_result['valeur_seuil'])
                        else:
                            seuils[api_indicator] = 0
                    except Exception as e:
                        print(f"Erreur seuil {api_indicator}: {e}")
                        seuils[api_indicator] = 0
            
            # Stocker les données de cette année
            comparison_data[str(annee)] = {
                'monthly_data': monthly_data,
                'seuils': seuils,
                'annee': annee,
                'atelier': atelier,
                'uap': uap
            }
        
        conn.close()
        
        return jsonify({
            'comparison_data': comparison_data,
            'atelier': atelier,
            'annee1': annee1,
            'annee2': annee2,
            'uap': uap,
            'aggregate': aggregate
        })
        
    except Exception as e:
        conn.close()
        print(f"Erreur dans get_charts_data_comparison: {e}")
        return jsonify({'error': f'Erreur serveur: {str(e)}'}), 500



import logging

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)
@app.route('/api/save_data', methods=['POST', 'OPTIONS'])
def save_data():
    # Handle preflight CORS requests
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'OK'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
        response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
        return response
    
    try:
        # Log incoming request
        logger.info(f"Received request: {request.method} {request.url}")
        logger.info(f"Content-Type: {request.content_type}")
        logger.info(f"Request data: {request.get_data()}")
        
        # Get JSON data
        if not request.is_json:
            logger.error("Request is not JSON")
            return jsonify({'error': 'Request must be JSON'}), 400
            
        data = request.get_json()
        logger.info(f"Parsed JSON data: {data}")
        
        if not data:
            logger.error("No data received")
            return jsonify({'error': 'No data received'}), 400

        # Validate and extract common fields
        mois = data.get('mois')
        annee = data.get('annee')
        atelier = data.get('atelier', '').strip()
        uap = data.get('uap', '').strip()

        # Enhanced validation
        errors = []
        
        if not mois or not isinstance(mois, int) or not (1 <= mois <= 12):
            errors.append("Le mois doit être un nombre entier entre 1 et 12")
            
        if not annee or not isinstance(annee, int) or annee <= 0:
            errors.append("L'année doit être un nombre entier positif")
            
        if not atelier:
            errors.append("L'atelier est requis")
            
        if not uap:
            errors.append("L'UAP est requis")
            
        if errors:
            logger.error(f"Validation errors: {errors}")
            return jsonify({'error': f"Erreurs de validation: {'; '.join(errors)}"}), 400

        # Connect to database
        try:
            conn = sqlite3.connect('qualite_ADI_final.sqlite')
            cursor = conn.cursor()
            logger.info("Database connection established")
        except Exception as db_error:
            logger.error(f"Database connection error: {db_error}")
            return jsonify({'error': f'Erreur de connexion à la base de données: {str(db_error)}'}), 500

        try:
            # Helper function to upsert into a table
            def upsert(table, value_field, value, value_type='INTEGER'):
                try:
                    logger.debug(f"Upserting into {table}: {value_field}={value}")
                    
                    # Check if record exists
                    cursor.execute(f"""
                        SELECT id FROM {table} 
                        WHERE mois=? AND annee=? AND atelier=? AND uap=?
                    """, (mois, annee, atelier, uap))
                    row = cursor.fetchone()

                    if row:
                        # Update existing record
                        logger.debug(f"Updating existing record with id={row[0]} in {table}")
                        cursor.execute(f"""
                            UPDATE {table} 
                            SET {value_field}=? 
                            WHERE id=?
                        """, (value, row[0]))
                    else:
                        # Insert new record
                        logger.debug(f"Inserting new record in {table}")
                        cursor.execute(f"""
                            INSERT INTO {table} (mois, annee, atelier, uap, {value_field}) 
                            VALUES (?, ?, ?, ?, ?)
                        """, (mois, annee, atelier, uap, value))
                        
                    return True
                    
                except Exception as table_error:
                    logger.error(f"Error in upsert for table {table}: {table_error}")
                    raise table_error

            # Process each data field
            operations = [
                ('pieces_retouchees', 'quantite', data.get('pieces_retouchees', 0), 'INTEGER'),
                ('pieces_rebutees', 'quantite', data.get('pieces_rebutees', 0), 'INTEGER'),
                ('pieces_controlees', 'quantite', data.get('pieces_controlees', 0), 'INTEGER'),
                ('cout_rebut', 'montant', data.get('cout_rebut', 0.0), 'REAL'),
                ('pieces_reclamees', 'quantite', data.get('pieces_reclamees', 0), 'INTEGER'),
                ('pieces_exportees', 'quantite', data.get('pieces_exportees', 0), 'INTEGER'),
                ('quantite_alertes', 'nombre_alertes', data.get('quantite_alertes', 0), 'INTEGER'),
                ('reclamations_officielles', 'nombre', data.get('reclamations_officielles', 0), 'INTEGER')
            ]

            # Execute all upserts
            for table, field, value, data_type in operations:
                # Convert and validate value
                if data_type == 'INTEGER':
                    try:
                        value = int(value) if value is not None else 0
                    except (ValueError, TypeError):
                        logger.warning(f"Invalid integer value for {field}: {value}, using 0")
                        value = 0
                elif data_type == 'REAL':
                    try:
                        value = float(value) if value is not None else 0.0
                    except (ValueError, TypeError):
                        logger.warning(f"Invalid float value for {field}: {value}, using 0.0")
                        value = 0.0
                
                # Ensure non-negative values
                if value < 0:
                    logger.warning(f"Negative value for {field}: {value}, using 0")
                    value = 0 if data_type == 'INTEGER' else 0.0
                
                upsert(table, field, value, data_type)

            # Commit transaction
            conn.commit()
            logger.info("All data saved successfully")
            
            # Success response
            response_data = {
                'message': 'Données enregistrées avec succès !',
                'status': 'success',
                'data_saved': {
                    'mois': mois,
                    'annee': annee,
                    'atelier': atelier,
                    'uap': uap
                }
            }
            
            response = jsonify(response_data)
            response.headers.add('Access-Control-Allow-Origin', '*')
            return response, 200

        except Exception as operation_error:
            logger.error(f"Error during database operations: {operation_error}")
            conn.rollback()
            return jsonify({
                'error': f'Erreur lors de l\'enregistrement: {str(operation_error)}',
                'status': 'error'
            }), 500

    except Exception as general_error:
        logger.error(f"General error in save_data: {general_error}")
        return jsonify({
            'error': f'Erreur interne du serveur: {str(general_error)}',
            'status': 'error'
        }), 500

    finally:
        try:
            if 'conn' in locals():
                conn.close()
                logger.debug("Database connection closed")
        except:
            pass

# Health check endpoint
@app.route('/api/health', methods=['GET'])
def health_check():
    try:
        # Test database connection
        conn = sqlite3.connect('qualite_ADI_final.sqlite')
        cursor = conn.cursor()
        cursor.execute("SELECT 1")
        conn.close()
        
        return jsonify({
            'status': 'healthy',
            'database': 'connected',
            'message': 'API is working correctly'
        }), 200
    except Exception as e:
        return jsonify({
            'status': 'unhealthy',
            'database': 'disconnected',
            'error': str(e)
        }), 500

# Error handlers
@app.errorhandler(404)
def not_found(error):
    return jsonify({
        'error': 'Endpoint not found',
        'message': 'The requested URL was not found on this server'
    }), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({
        'error': 'Internal server error',
        'message': 'An unexpected error occurred'
    }), 500




@app.route('/quality_manager/graphique')
@login_required
def quality_manager_graphique():
    return render_template('quality_manager/graphique.html')

@app.route('/viewer/dashboard')
def viewer_dashboard():
    # Vous pouvez ajouter ici toute logique nécessaire
    # Par exemple, vérifier l'authentification, charger des données, etc.
    return render_template('viewer/dashboard.html')


@app.route('/quality_manager/dashboard')
def graphique():
    return render_template('quality_manager/dashboard.html')


@app.route('/api/cnq_ytd_value', methods=['GET'])
def get_cnq_ytd_value():
    """Récupère la valeur YTD CNQ depuis la colonne valeur_ytd de cnq"""
    try:
        year = request.args.get('annee', default=datetime.now().year, type=int)
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Récupérer la valeur YTD (non-zéro) du mois 1 de l'année
        cursor.execute(
            '''SELECT valeur_ytd FROM cnq 
               WHERE annee = ? AND mois = 1 AND atelier = "Total" AND valeur_ytd > 0
               ORDER BY created_at DESC LIMIT 1''',
            (year,)
        )
        result = cursor.fetchone()
        conn.close()
        
        if result and result[0]:
            return jsonify({
                'valeur_ytd': float(result[0]),
                'annee': year,
                'status': 'success'
            })
        else:
            return jsonify({
                'valeur_ytd': 0,
                'annee': year,
                'status': 'no_data'
            })
    except Exception as e:
        print(f"❌ Erreur get_cnq_ytd_value: {str(e)}")
        return jsonify({
            'error': str(e),
            'status': 'error',
            'valeur_ytd': 0
        }), 500


@app.route('/admin/update_cnq_ytd', methods=['POST'])
def update_cnq_ytd():
    """Met à jour la valeur YTD du CNQ dans la colonne valeur_ytd de cnq"""
    try:
        cnq_ytd_value = float(request.form.get('cnq_ytd_value'))
        cnq_ytd_year = int(request.form.get('cnq_ytd_year'))
        
        # Validation
        if cnq_ytd_value < 0 or cnq_ytd_value > 100:
            flash('La valeur YTD CNQ doit être entre 0 et 100%', 'error')
            return redirect(url_for('admin_thresholds'))
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Vérifier si une entrée existe déjà pour mois=1 et atelier="Total"
        cursor.execute(
            '''SELECT id FROM cnq 
               WHERE annee = ? AND mois = 1 AND atelier = "Total"''',
            (cnq_ytd_year,)
        )
        existing = cursor.fetchone()
        
        if existing:
            # Mettre à jour - seulement la colonne valeur_ytd
            cursor.execute(
                '''UPDATE cnq SET valeur_ytd = ? 
                   WHERE annee = ? AND mois = 1 AND atelier = "Total"''',
                (cnq_ytd_value, cnq_ytd_year)
            )
            message = f'Valeur YTD CNQ {cnq_ytd_year} mise à jour avec succès'
        else:
            # Insérer une nouvelle entrée avec mois=1
            cursor.execute(
                '''INSERT INTO cnq (mois, annee, atelier, valeur, seuil, valeur_ytd) 
                   VALUES (1, ?, "Total", 0, 0, ?)''',
                (cnq_ytd_year, cnq_ytd_value)
            )
            message = f'Valeur YTD CNQ {cnq_ytd_year} ajoutée avec succès'
        
        conn.commit()
        conn.close()
        
        flash(message, 'success')
        print(f"✅ CNQ YTD mise à jour dans table cnq: {cnq_ytd_year} = {cnq_ytd_value}%")
        
    except ValueError as e:
        flash(f'Erreur de validation : {str(e)}', 'error')
        print(f"❌ Erreur validation CNQ YTD: {str(e)}")
    except Exception as e:
        flash(f'Erreur : {str(e)}', 'error')
        print(f"❌ Erreur CNQ YTD: {str(e)}")
    
    return redirect(url_for('admin_thresholds'))

# À ajouter dans votre app.py

def migrate_cnq_table():
    """Ajoute la colonne valeur_ytd si elle n'existe pas"""
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Vérifier si la colonne existe déjà dans cnq
        cursor.execute("PRAGMA table_info(cnq)")
        columns = [column[1] for column in cursor.fetchall()]
        
        if 'valeur_ytd' not in columns:
            print("⚠️  Ajout de la colonne valeur_ytd à la table cnq...")
            cursor.execute(
                'ALTER TABLE cnq ADD COLUMN valeur_ytd DECIMAL(10, 3) DEFAULT 0'
            )
            conn.commit()
            print("✅ Colonne valeur_ytd ajoutée à cnq avec succès")
        else:
            print("ℹ️  La colonne valeur_ytd existe déjà dans cnq")
        
        conn.close()
    except Exception as e:
        print(f"❌ Erreur lors de la migration: {str(e)}")

        
if __name__ == '__main__':
    # Initialiser la base de données
    init_database()
    migrate_cnq_table()  # Ajoute la colonne si manquante
    print("Base de données initialisée avec succès!")
    print("Utilisateur admin créé:")
    #print("  - Nom de compte: skander.chniti")
    #print("  - Mot de passe: @@skanderchniti1$")
    # print("  - Rôle: Administration")
    print("\nServeur démarré sur http://localhost:5000")
    
    # Lancer en mode debug localement
    app.run(debug=True, host='0.0.0.0', port=5000)